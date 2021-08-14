import os
import logging
import xmltodict
import shutil
import zipfile
import pathlib
import re
import natsort

from typing import OrderedDict
from pptx import Presentation
from functools import reduce
from lxml import etree
# from zipfile import ZipFile
# from shutil import copyfile, rmtree, move, make_archive


logger = logging.getLogger (__name__)

fq_empty = "mob/beagle/test_resources/Empty.pptx"


def base_name (fqfn) :
    """
    returns file name, i.e. fqfn
    """
    return reduce(lambda x,f : f(x),
                  [fqfn , 
                   os.path.basename, 
                   os.path.splitext]) [0]


def file_to_work_dir (fqfn) :
    """
    returns /tmp/{fqfn}
    """
    l = ["tmp", base_name (fqfn)] # removed '/'
    return os.path.join (*l) 


def token_to_work_dir (ft) :
    """
    """
    l = ["tmp", ft] # removed '/'
    return os.path.join (*l) 


def token_is_open (ft) :
    """
    """
    fp = token_to_work_dir (ft)
    if not os.path.exists (fp) :
        raise ValueError (f"{ft} is not open")
    return fp


def file_is_open (fn) :
    """
    """
    fp = file_to_work_dir (fn)
    if not os.path.exists (fp) :
        raise ValueError (f"{fp} is not open")
    return fp


def file_is_not_open (fn) :
    """
    """
    fp = file_to_work_dir (fn)
    if os.path.exists (fp) :
        raise ValueError (f"{fn} is already open")
    return fp


def unpack (fqfn) :
    """
    **fqfn should be the full system path
    unpack('C:/Users/RichaM/Documents/code/parse_task/tmp/Onboarding.pptx')

    # fqfn = '/tmp/Onboarding' (error)
    removed '/' from "file_to_work_dir"

    fqfn: input deck name

    fqfn : Onboarding
    fp = /tmp/Onboarding 

    fqfn : f"{base_path}/{presentation}/Onboarding.pptx"

    /mnt/input/{messsage-fn}
    """
    fp = file_is_not_open (fqfn) # /tmp/{fqfn}
    # fp = '/tmp/new_file.pptx'
    with zipfile.ZipFile(fqfn) as z:
        z.extractall(fp)
    return fp


def pack (ft, *args, **kwargs) :
    """
    """
    fp = token_is_open (ft)
    d = kwargs.get ("file", ".".join ([ft, "pptx"]))
    z = shutil.make_archive (bn , "zip", file_to_work_dir (fqfn))
    shutil.move (z, d)
    logger.debug (f"wrote {d}")
    return d


def close (ft) :
    """
    """
    fp = token_is_open (ft)
    shutil.rmtree (fp)
    logger.debug (f"deleted {ft}")
    return


def new(ft):
    """
    create, move and unzip the empty output deck

    "/tmp/x" -> /tmp/x.pptx
    """
    fq_empty = f"/tmp/{path}"
    # create
    prs = Presentation()
    prs.save(fq_empty)
    # unzip
    fp = unpack (fq_empty)
        
    return fp


def get_fld_f_names (asset):
    """
    generates folder and file names
    """
    sp = asset.split ('/')
    if '_rels' in asset: # slides/_rels/slide2.xml.rels
        f = sp[-1]
        fld = f'{sp[-3]}/{sp[-2]}'
    else:
        fld, f = sp[-2], sp[-1]

    return fld, f


def xml_to_dict (path):
    """
    convert xml file to dict
    """
    with open(path) as xml_file:
        data_dict = xmltodict.parse(xml_file.read())
        xml_file.close()
    if isinstance(data_dict["Relationships"]["Relationship"], list):
        data = sorted(data_dict["Relationships"]["Relationship"], key=lambda item: int(item['@Id'].split('Id')[1]))
    else:
        data = [data_dict["Relationships"]["Relationship"]]
    return data


def write_output_xml(tree, f):
    """
    creates an XML file using the tree
    """
    tree.write(f, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    
    return f


def modify_output_xml(tree, relation, rel_tag):
    """
    update existing xml file
    """
    prev = tree.find(rel_tag)
    prev.addnext(relation)
    
    return


def build_tree (path):
    """
    pass the path of the xml file to enable the parsing process
    returns root and tree object of the xml file
    """
    parser = etree.XMLParser(remove_blank_text=True)
    tree = etree.parse(path, parser)
    # root = tree.getroot()    
    return tree


def refactor_assets_and_rels(k, v, ctx):
    """
    generates refactored name for an asset,
    update the refactored count and
    creates assets in the output deck
    """
    lasts = ctx["lasts"]
    refactored_fns = ctx["refactored_fns"]

    # ['tmp', 'Onboarding', 'ppt', 'media', 'image10.png']
    # ['tmp', 'Onboarding', 'ppt', 'slideLayouts', '_rels', 'slideLayout2.xml.rels']
    l = k.split("/")
    fn_stub = re.findall(r'(\w+?)(\d+)', k)[0][0]

    if l[-2] is "_rels":
        file_ext = ".xml.rels"
        lasts_key = f"{l[-3]}/_rels/{fn_stub}"
    else:
        file_ext = l[-1].split(".")[-1]
        lasts_key = f"{l[-2]}/{fn_stub}"
    
    new_index, lasts = next_index_for(lasts_key, lasts)
    l[-1] = f"{fn_stub}{new_index}.{file_ext}"
    
    # _rels/slideLayout1.xml.rels
    # slideLayouts/slideLayout1.xml
    # media/image5.png
    new_name = '/'.join(l[-2:])
    
    refactored_fns = {**refactored_fns, k: new_name}

    ctx ["lasts"] = lasts
    ctx ["refactored_fns"] = refactored_fns
    
    return ctx


def remove_default_assets(ctx, otags):
    """
    remove duplicate tags from ctx['pxr_relations']
    """
    pxr_relations = ctx['pxr_relations']
    for k in pxr_relations.keys():
        if '/' not in k and k in otags:
            del pxr_relations[k]
    
    ctx['pxr_relations'] = {**pxr_relations}
    return ctx
            

def refactor_relations(ctx):
    """
    """
    pxr_relations = ctx['pxr_relations']
    refactored_fns = ctx['refactored_fns']
    
    for k,v in pxr_relations.items():
        if 'ppt' in k:
            name = refactored_fns[k].split('ppt')[-1]
        else:
            name = '../' + '/'.join (refactored_fns[k].split('/')[-2:])
        v.set('Target', name)
    return ctx
    

def refactor_rIds(ctx):
    """
    set refactored rId
    """
    lasts = ctx['lasts']
    pxr_relations = ctx['pxr_relations']
    max_rId = lasts['rId']
    
    for k,v in pxr_relations.items():
        max_rId += 1
        v.set ('Id', f"rId{str(max_rId)}")
    
    
      

def remove_default_files(ft):
    """
    remove all the default assets
    return maximum rId
    """
    default = ['slideMasters', 'slideLayouts', 'theme']
    
    for root, dirs, files in os.walk(ft):
        for f in files:
            if 'slideMasters' in f or 'slideLayouts' in f or 'theme' in f:
                os.remove(os.path.join(root, f))

    
    return


def assetfn_to_relfn(ft, fqfn):
    """
    {ft}/customXml/item1.xml -> {ft}/customXml/_rels/item1.xml.rels
    {ft}/ppt/slides/slide2.xml  ->  {ft}/ppt/slides/_rels/slide2.xml.rels
    """
    # if '/' in fqfn:
    l = fqfn.split("/")

    if l[-3] is not "ppt":
        r = f"{ft}/{l[-2]}/_rels/{l[-1]}.rels"
    else:
        r = f"{ft}/ppt/{l[-2]}/_rels/{l[-1]}.rels"
    
    return r


def relfn_to_assetfn(ft, fqfn):
    """
    {ft}/customXml/_rels/item1.xml.rels -> {ft}/customXml/item1.xml
    {ft}/ppt/slides/_rels/slide2.xml.rels -> {ft}/ppt/slides/slide2.xml
    """
    # if '/' in fqfn:
    l = fqfn.split("/")
    l.pop(-2)
    l[-1] = l[-1][:-5]
    r = f"{ft}/{l[-2]}/{l[-1]}"
    
    return r


def short_assetfn_to_long_assetfn(ft, short):
    """
    ../customXml/item1.xml -> customXml/item1.xml
    slides/slide2.xml  -> /ppt/slides/slide2.xml
    ../slideLayout2.xml -> /ppt/slideLayout2.xml
    """
    if '..' in short:
        if os.path.isfile(f"{ft}/{short[:3]}"):
            return short[:3]
        elif os.path.isfile(f"{ft}/ppt/{short[:3]}"):
            return f"ppt/{short[:3]}"
    
    r = f"ppt/{short}"
    
    return r


def is_manadatory(i):
    """
    return true if asset should be copied
    """
    attrib = i.attrib['Target']
    return '/' in attrib and \
        not attrib.startswith("slides") and \
            not attrib.startswith("slideLayouts") and \
                not attrib.startswith("slideMasters") and \
                    not attrib.startswith("theme")


def walk_asset_tree(ft, asset, ctx):
    
    rf = assetfn_to_relfn(ft, asset)

    # assetfn = f"{ft}/ppt/{asset}"
    assetfn = short_assetfn_to_long_assetfn(ft, asset)
    assets = ctx["assets"]
    rels = ctx["rels"]
    assets = {**assets, assetfn: build_tree(assetfn)}

    if os.path.exists(rf):
        rels_contents = build_tree(rf)
        rels = {**rels, rf: rels_contents}
        for r in rels_contents.getroot():
            t = r.attrib['Target'][:3] # ../media/image1.png
            ctx = walk_asset_tree(ft, t, ctx)

    ctx["assets"] = assets
    ctx["rels"] = rels

    return ctx

# /mnt/input
# /mnt/output


def next_index_for(type, lasts):
    """
    
    """
    current = lasts.get(type, 0)
    current = current + 1
    lasts[type] = current
    return current, lasts


def refactor_content(k, v, ctx):
    """
    refactoring content of the rel files
    and then saving it
    
    refactored_fns = {'/tmp/Onboarding/ppt/media/image10.png': 'media/image5.png'}
    """
    refactored_fns = ctx['refactored_fns']
    ft = '/'.join(k.split('/')[:-3])
    
    for relation in v.getroot():
        target = relation.attrib['Target']
        # {ft}/media/image10.png  -> media/image5.png
        if '..' in target:
            new_name = f"../{refactored_fns[f'{ft}/{target[3:]}']}"
        else:
            new_name = refactored_fns[f"{ft}/{target}"]
        
        relation.set ('Target', new_name)
    
    return


def build_required_assets(ft, ctx):
    pxr = f'{ft}/ppt/_rels/presentation.xml.rels'
    tree = build_tree (pxr)
    root = tree.getroot()
    d = filter(is_manadatory, root)
    
    for asset in d:
        ctx = walk_asset_tree(ft, asset, ctx)
    return ctx


def build_mandatory_assets(ft, ctx):
    """
    """
    mand_assets_list = ['slideMasters', 'slideLayouts', 'theme']
    
    mandatory_assets = ctx['mandatory_assets']
    
    for fl in mand_assets_list:
        for root, dirs, files in os.walk (f'{ft}/ppt/{fl}'):
            # for windows
            if '\\' in root:
                root = root.replace('\\', '/')
            
            f_list = natsort.natsorted (files)
            for f in f_list:
                asset = f"{root}/{f}"
                mandatory_assets = {**mandatory_assets, asset: build_tree (asset)}
    
    ctx['mandatory_assets'] = mandatory_assets
    
    return ctx


#def build_assets(ft, pxr, assets, slide=None):
def build_assets_and_rels(s, ft, ctx):
    """
    assets = {fq-destination, file}
    req:
     assets (with all the files)
    """
    s = f'slides/slide{str(s)}.xml'
    ctx = walk_asset_tree(ft, s, ctx)

    # if 'placeholder':
    #     pass
    # else:
    #     tree = build_tree (pxr)
    #     root = tree.getroot()
    #     for relation in root:
    #         attrib = relation.attrib
    #         if 'http' not in attrib['Target']:
    #             if '../' in attrib['Target']:
    #                 tar = f'{ft}/{attrib[3:]}'
    #                 assets[tar] = build_tree (tar)

    #                 if 'xml' in attrib['Target']:
    #                     fld = attrib['Target'].split('/')[-2]
    #                     fl = attrib['Target'].split('/')[-1]
    #                     pxr = f'{ft}/ppt/{fld}/_rels/{fl}'
                
    #                     if os.path.exists (pxr):
    #                         build_assets (ft, pxr, assets)
                
    #             else:
    #                 new_tar = pxr.split('/')[-3]
    #                 tar = f"{ft}/ppt/{new_tar}/{attrib['Target']}"
    #                 assets[tar] = build_tree (tar)
    
    return ctx


# def build_rels(ft, assets, rels):
#     """
#     add a check if the asset is having a rel file
#     if yes then add it
    
#     considering 'ft' as full path, "/tmp/{input_deck}
#     """
#     for i in assets.keys():
#         f = assetfn_to_relfn(ft, i)         
#         if os.path.isfile(f):
#             rels[f] = build_tree (f)
    
#     return rels


def build_content_types(ft, ctx):
    """
    add content_type of new assets in [Contant_Types].xml file
    
    considering 'ft' as full system input deck path 
    """

    f = '[Content_Types].xml'
    inp_path = '/'.join ([ft, f])

    content_types = ctx["content_types"]

    tree1 = build_tree (inp_path)

    for relation in tree1.getroot():
        if 'Override' in relation.tag:
            attrib = relation.attrib['PartName'][1:] # /ppt/slides/slide1.xml
            try:
                cnt = attrib.split ('ppt/')[-1]
                ini = f'{ft}/ppt'
            except:
                cnt = attrib
                ini = f'{ft}'
            
            name = f'{ini}/{cnt}'
            
            if name in ctx["assets"].keys() or name in ctx["rels"].keys() or name in ctx['mandatory_assets']:
                content_types = {**content_types, name: relation}
        else:
            attrib = relation.attrib['Extension']
            content_types = {**content_types, attrib: relation}

    ctx["content_types"] = content_types
    return ctx


def build_properties(ft, ctx):
    
    inp_path = '/'.join([ft, 'ppt'])
    properties = ctx["properties"]
    
    for i in os.listdir(inp_path):
        f = f'{inp_path}/{i}'
        if os.path.isfile(f):
            properties = {**properties, f: build_tree (f)}
    
    ctx["properties"] = properties

    return ctx


def build_pxr_file(ft, ss, ctx):
    """
    """
    refactored_fns = ctx['refactored_fns']
    lasts = ctx['lasts']
    
    pxr = f"{ft}/ppt/_rels/presentation.xml.rels"
    
    pxr_relations = ctx["pxr_relations"]
    ss = [f"slides/slide{s}.xml" for s in ss]
    
    root = build_tree(pxr).getroot()
    
    for relation in root:
        target = relation.attrib['Target']
        if 'slides/slide' not in target:
            pxr_relations = {**pxr_relations, f'{ft}/{target}': relation}
        elif target in ss:
            pxr_relations = {**pxr_relations, f'{ft}/{target}': relation}
    
    ctx["pxr_relations"] = pxr_relations
    
    return ctx


def apply_assets(ft, ctx):
    """
    creates all assets in the output deck 
    """    
    for k,v in ctx['assets'].items():
        # change the asset key to the next index per asset type
        #       old             ->              new
        # .../ppt/slides/slide2.xml ->  .../ppt/slides/slide1.xml
        # .../customXml/items2.xml  ->  .../customXml/items1.xml
        ctx = refactor_assets_and_rels(k, v, ctx)
    
    return ctx
    

def apply_rels(ft, k, v, ctx):
    """
    creates all rel files of assets in the output deck 
    """
    for k,v in ctx['rels'].items():
        # change the asset key to the next index per asset type
        #       old             ->              new
        # .../ppt/slides/slide2.xml ->  .../ppt/slides/slide1.xml
        # .../customXml/items2.xml  ->  .../customXml/items1.xml
        refactor_content(k, v, ctx)
        ctx = refactor_assets_and_rels(k, v, ctx)
    
    return ctx
 
    # fld, f = get_fld_f_names(k)
    # assets_nm, assets_cnt = {}, {}

    # ext = ''.join(pathlib.Path(f).suffixes) # .xml (or) .xml.rels
    # name = re.findall(r'(\w+?)(\d+)', f)[0][0] # slide, slideLayout, theme, slideMaster

    # if f'{fld}/{name}' in assets_cnt.keys():
    #     cnt = assets_cnt[f'{fld}/{name}']+1
    # else:
    #     cnt = 1
    
    # new_name = f'{name}{cnt}{ext}'
    # tree = v
    
    # if 'ppt' in k:
    #     write_output_xml(tree, f"{ft}/ppt/{fld}/{new_name}")
    # else:
    #     write_output_xml(tree, f"{ft}/{fld}/{new_name}")

    # assets_nm[f'{fld}/{f}'] = f'{fld}/{new_name}'
    # assets_cnt[f'{fld}/{name}'] = cnt
    
    # return assets_nm, assets_cnt


def apply_mandatory_assets(ft, ctx):
    """
    """
    for k, v in ctx['mandatory_asset'].items():
        ctx = refactor_assets_and_rels(k, v, ctx)
    
    return ctx


def apply_content_types(ft, ctx):
    """
    need to perform refctoring of names before making changes in the output deck's file
    """
    con = '[Content_Types].xml'
    f = '/'.join ([ft, con])
    tree = build_tree (f)
    content_types = ctx['content_types']
    refactored_fns = ctx['refactored_fns']
    
    for k,v in content_types.items():
        if k in refactored_fns.keys():
            # slideMasters/slideMaster1.xml -> /ppt/slideMasters/slideMaster1.xml
            if 'ppt' in k:
                nm = refactored_fns[k]
                new_name = f"/ppt/{nm}"
            else:
                # customXml/itemProps2.xml -> /customXml/itemProps2.xml
                new_name = f"/{nm}"
            
            v.set('PartName', new_name)
            modify_output_xml(tree, v, v.tag)
        else:
            # png
            # /ppt/presentation.xml
            modify_output_xml(tree, v, v.tag)
    
    write_output_xml(tree, f)
    
    return ctx


def apply_properties(ft, properties):
    """
    considering 'ft' as full system output deck's path
    """
    
    mergables = ['commentAuthors.xml', 'tableStyles.xml']
    sing_prop = ['viewProps.xml', 'presProps.xml']
    ignore = ['revisionInfo.xml']
    
    for k,v in properties.items():
        i = k.split('/')[-1]
        f = f"{ft}/ppt/{i}"
        
        if os.path.isfile (f):
            tree1 = v
            root1 = tree1.getroot()
            tree2 = build_tree (f)
            root2 = tree2.getroot()
            
            etag = f"{root1[0].tag}"
            if i in mergables:
                try:
                    for relation in [etag]:
                        for ele in root1.findall (relation):
                            root2.append(ele)
                except IndexError:
                    logger.debug(f"list index {i} out of range")
            elif i in sing_prop:
                if i == 'presProps.xml':
                    uris = {}
                    out_lis = []
                    nm = root1.nsmap['p']
                    ext_tag = f"{{{nm}}}extLst"
                    
                    for relation in [etag]:
                        
                        fp = root1.find (ext_tag)
                        for ele in fp:
                            attrib = ele.attrib
                            if attrib['uri'] not in uris.keys():
                                uris[attrib['uri']] = ele

                    for relation in [f"{root2[0].tag}"]:
                        fp = root2.find (ext_tag)
                        for ele in fp:
                            attrib = ele.attrib
                            out_lis.append (attrib['uri'])
                    
                    for k,v in uris.items():
                        if k not in out_lis:
                            modify_output_xml(tree2, v, ext_tag)
                            # tag1 = tree2.find(ext_tag)
                            # tag1.append(v)

            write_output_xml(tree2, f)
        else:
            write_output_xml(tree1, f)
    
    return properties


def apply_pres_files(ft, ctx):
    """
    """
    pxr = f"{ft}/ppt/_rels/presentation.xml.rels"
    
    otags = []
    
    root = build_tree(pxr).getroot()
    for relation in root:
        otags.append(relation.attrib['Target'])
    
    ctx = remove_default_assets(ctx, otags)
    ctx = refactor_relations(ctx)
    
    return ctx


def process_message(msg, ctx):
    logger.debug(f"processing message {msg}")
    input_deck = f"/mnt/input/{msg['d']}"
    
    ft = unpack (input_deck)
    
    # ft: extracted input_deck dir (/tmp/{input_deck})
    
    
    ss = msg.get("s", None)
    if not ss:
        prs = Presentation(ft)
        ss = list(range(1, len(prs.slides._sldIdLst) + 1))

    ctx = build_required_assets(ft, ctx)
    ctx = build_mandatory_assets(ft, ctx)

    for s in ss:
        ctx = build_assets_and_rels(s, ft, ctx)
        #rels = build_rels(ft, assets, rels, lasts)
    
    ctx = build_content_types(ft, ctx)
    ctx = build_properties(ft, ctx)
    ctx = build_pxr_file(ft, ss, ctx)

    return ctx


def write_output(ft, ctx, output_deck):
    
    ctx = refactor_assets_and_rels(ctx)

    ctx = write_assets_and_rels(ft, ctx)
    ref_nm, ref_cnt = apply_assets(ft, ctx["assets"], ctx['refactored_cnt'])
    ctx['refactored_fns'] = {**ref_nm}
    ctx['refactored_count'] = {**ref_nm}
    
    ref_nm, ref_cnt = apply_rels(ft, ctx["rels"], ctx['refactored_name'], ctx['refactored_cnt'])
    ctx['refactored_fns'] = {**ref_nm}
    ctx['refactored_count'] = {**ref_nm}
    
    apply_content_types(ft, ctx)
    apply_properties(ft, ctx["properties"])
    apply_pres_files(ft, ctx)

    pack(ft, file=output_deck)
    close(ft)
    return output_deck


def render_deck_effect(ctx):
    """
    process msgs and write output deck to fqfn

    fqfn should be /mnt/output/{render_id}
    """
    msgs = ctx.get ("msgs", [])
    if not msgs:
        logger.error("No messages in context")
    else:
        render_id = msgs.pop(0)
        output_deck = f"/mnt/output/{render_id}.pptx"
        
        # ft: full system path of empty output deck
        ft = new(output_deck)
        
        remove_default_files(ft)
        lasts = {'rId': 6}
        
        ctx = reduce(process_message, msgs, {"lasts": lasts,
            "refactored_fns": OrderedDict(),
            "assets": OrderedDict(),
            "mandatory_assets": OrderedDict(),
            "rels": OrderedDict(),
            "content_types": OrderedDict(),
            "properties": OrderedDict(),
            "pxr_relations": OrderedDict()})

        
        logger.debug(f"writing output deck to {output_deck}")
        write_output(ft, ctx, output_deck, lasts)

    return ctx

# continue
