import os
import shutil
import json
import xmltodict
import natsort
import zipfile
import re
import pathlib

from pptx import Presentation
from lxml import etree
from collections import OrderedDict
from functools import reduce


def unpack (src, des):
    """
    unzip the src location deck to des location
    """
    with zipfile.ZipFile(src, 'r') as zip_ref:
        zip_ref.extractall(des)
    return des


def pack (path, file_name):
    """
    zip extracted deck to get output deck
    """
    length = len(path)
    zipf = zipfile.ZipFile('output/'+f'Test_{file_name}.pptx', 'w', zipfile.ZIP_DEFLATED)
    for root, dirs, files in os.walk(path):
        folder = root[length:] # path without "parent"
        for file in files:
            zipf.write(os.path.join(root, file), os.path.join(folder, file))
    zipf.close()
    return


    """
    returns the base path
    """
    return os.getcwd()


def build_tree (path):
    """
    pass the path of the xml file to enable the parsing process
    """
    parser = etree.XMLParser(remove_blank_text=True)
    tree = etree.parse(path, parser)
    root = tree.getroot()    
    return root, tree


def max_rId (emp):
    """
    find the largest rId
    """
    path = f'{emp}/ppt/_rels/presentation.xml.rels'
    root, tree = build_tree(path)

    rIds = []

    for relation in root:
        attrib = relation.attrib
        rId = int(attrib.get('Id').split('Id')[-1])
        rIds.append(rId)
    return {'rId': max(rIds)}


def xml_to_dict (path):
    """
    convert xml file to dict
    """
    with open(path) as xml_file:
        data_dict = xmltodict.parse(xml_file.read())
        xml_file.close()
    if isinstance(data_dict["Relationships"]["Relationship"], list):
        data = sorted(data_dict["Relationships"]["Relationship"], key=lambda item: int(item['@Id'].split('Id')[1]))
    else:
        data = [data_dict["Relationships"]["Relationship"]]
    return data


def total_slides (path):
    """
    calculate total number of slides
    """
    prs = Presentation(path)
    tot_slides = len(prs.slides._sldIdLst)
    return tot_slides


def first_slide (path):
    """
    calculate the rId of the first slide
    """
    root, _ = build_tree (path)
    slide = 'slide1.xml'
    
    for relation in root:
        attrib = relation.attrib
        if slide in attrib['Target']:
            return int(attrib['Id'].split('Id')[-1])
    
    return


def build_rels (assets_name):
    """
    list latest .rels files
    """
    l = [i for i in assets_name.values()]
    rel = natsort.natsorted ([i for i in l if '_rels' in i])
    
    return rel


def build_slides (pp):
    """
    generate a list of slides
    """
    ts = total_slides (pp) + 1
    sld = list(range(1, ts))
    
    return sld


def repair_path (assets_name):
    """
    replacing '\\' with '/'
    """
    dup_dict = dict(assets_name)
    for k,v in dup_dict.items():
        if '\\' in k:
            key = k.replace('\\', '/')
            val = v.replace('\\', '/')
            del assets_name[k]
            dict_1[key] = val
    return assets_name


def remove_dup (itags, assest_rels, otags):
    """
    remove duplicate relationships from the rels files
    """
    l1 = itags[:]
    
    for i in l1:
        if '/' not in i:
            if i in otags:
                itags.remove(i)
                del assest_rels[i]
    
    return itags   


def new (path):
    """
    create, move and unzip the empty output deck
    """
    fq_empty = "resources/Empty.pptx"
    
    # create
    prs = Presentation()
    prs.save (fq_empty)
    
    # move
    d = ".".join ([path, "pptx"]) # "output/41.pptx"
    shutil.move (fq_empty, d)
    
    # unzip
    unpack (d, d.split('.')[0])
    
    os.remove (d)
    
    return d.split ('.')[0]


def make_dir (ip, ft, input_deck):
    """
    create input deck directories which does not exists in the output deck
    """
    for i in os.walk (ip):
        fld = i[0].split (input_deck)[-1]
        if fld:
            loc = f"{ft}{fld}"
            if not os.path.exists (f'{ft}/{fld}'):
                os.makedirs (f'{ft}/{fld}')
    
    return ft


def make_structure (ip, ft, input_deck):
    """
    create structure of first input deck
    """
    for i in os.walk (ip):
        fld = i[0].split (input_deck)[-1]
        if fld:
            loc = f"{ft}{fld}"
            if 'ppt' not in loc and (input_deck not in loc):
                shutil.rmtree (f'{ft}/{fld}')
                shutil.copytree (f'{ip}/{i[0].split (input_deck)[-1]}', f'{ft}/{fld}')
    
    return ft


def build_assets (pxr, input_deck, assets, ss=None):
    """
    list the assets that needs to be modified in output deck
    """
    data = xml_to_dict (pxr)
    # tmp_path = "/Users/sanskar/Desktop/Development/MS Bing/Dev-zip/tmp/41"
    tmp_path = "C:/Users/RichaM/Documents/code/parse_task/tmp/41"
    
    if ss:
        # presentations = "/Users/sanskar/Desktop/Development/MS Bing/Dev-zip/presentations"
        presentations = "C:/Users/RichaM/Documents/code/parse_task/presentations"
        tot_slides = total_slides (f'{presentations}/{input_deck}.pptx')
        first_slide_id = first_slide (pxr)

        files = []
        for i in data:
            curr_rId = int(i['@Id'].split ('Id')[1])
            if (first_slide_id > curr_rId) or (curr_rId > (first_slide_id + tot_slides - 1)):
                if 'slideLayouts' not in i['@Target'] and 'slideMasters' not in i['@Target'] and 'theme' not in i['@Target']:
                    files.append (i['@Target'])
        
        for i in files:
            if '/' in i:
                a = i.split ('/')
                fld, fl = a[0], a[1]
                if os.path.exists (f'{tmp_path}/{input_deck}/ppt/{fld}/_rels'):
                    if os.path.isfile (f'{tmp_path}/{input_deck}/ppt/{fld}/_rels/{fl}.rels'):
                        assets.append (f'{fld}/_rels/{fl}.rels')
        
        assets = assets + files
        
        for s in ss:
            slide = f'slide{str(s)}.xml'
            assets.append ([i["@Target"] for i in data if slide in i["@Target"] and "http" not in i["@Target"]][0])
            assets.append (f'slides/_rels/{slide}.rels')
            build_assets (f'{tmp_path}/{input_deck}/ppt/slides/_rels/{slide}.rels', input_deck, assets)

    else:
        for i in data:
            # handling duplicacy
            if i["@Target"] in assets or i["@Target"][3:] in assets:
                pass
            
            elif "http" not in i["@Target"]:    
                if '../' in i["@Target"]:
                    assets.append (i['@Target'][3:])
                else:
                    new_tar = pxr.split ('/')[-3]
                    assets.append(f'{new_tar}/{i["@Target"]}')

                if ".." in i['@Target'] and "xml" in i['@Target']:
                    fld = i['@Target'].split ('..')[1].split ('/')[1]
                    fl = i['@Target'].split ('..')[1].split ('/')[2]
                    pxr = f"{tmp_path}/{input_deck}/ppt/{fld}/_rels/{fl}.rels"

                    if os.path.exists (pxr):
                        # handling rels files
                        assets.append (pxr.split ('ppt/')[1])
                        build_assets (pxr, input_deck, assets)

    return assets


def build_names (asset):
    """
    generates folder and file names
    """
    if '_rels' in asset: # slides/_rels/slide2.xml.rels
        sp = asset.split ('/')
        f = sp[-1]
        fld = f'{sp[0]}/{sp[1]}'
    elif '../' in asset:
        _, fld, f = asset.split ('/')
    else:
        fld, f = asset.split ('/')

    return fld, f


def rename (ft, ip, fld, f, assets_cnt): # fld=media, fl=image1.png
    """
    rename an asset
    """
    d1 = OrderedDict()

    ext = ''.join(pathlib.Path(f).suffixes)
    name = re.findall(r'(\w+?)(\d+)', f)[0][0]
    
    if os.path.exists (f'{ip}/ppt/{fld}/{f}'):
        path = f'{ip}/ppt'
    else:
        path = ip
    
    if f'{fld}/{name}' in assets_cnt.keys():
        count = assets_cnt[f'{fld}/{name}']+1
    else:
        count = 1
    new_name = f'{name}{count}{ext}'
    if 'ppt' in path:
        shutil.copy(f'{path}/{fld}/{f}', f"{ft}/ppt/{fld}/{new_name}")
    else:
        shutil.copy(f'{path}/{fld}/{f}', f"{ft}/{fld}/{new_name}")
    d1[f'{fld}/{f}'] = f'{fld}/{new_name}'
    assets_cnt[f'{fld}/{name}'] = count
    return d1


def apply_mend_assets (ip, ft, d, assets_name, assets_cnt):
    """
    copy mandatory assests
    """
    nm = dict(assets_name)
    cnt = dict(assets_cnt)
    fl = ['slideLayouts', 'theme', 'slideMasters']
    des = f'{ft}/ppt'
    if d == 1:
        for f in fl:
            count = 0
            for i in os.walk (f'{ip}/ppt/{f}'):
                count = len (i[2])
            if os.path.exists (f'{des}/{f}'):
                shutil.rmtree (f'{des}/{f}')
            shutil.copytree (f'{ip}/ppt/{f}', f'{des}/{f}')
            cnt [f] = count
            if os.path.exists (f'{ip}/ppt/{f}/_rels'):
                cnt [f'{f}/_rels'] = count
            else:
                cnt [f'{f}/_rels'] = count

            for i in os.walk (f'{ip}/ppt/{f}'):
                fld = i[0].split ('ppt/')[1]
                if '\\' in fld:
                    fld = fld.replace ('\\', '/')
                fl_list = natsort.natsorted (i[2])
                for j in fl_list:
                    nm [f'{fld}/{j}'] = f'{fld}/{j}'
    else:
        for f in fl:
            if os.path.exists (f'{des}/{f}'):
                for j in os.walk (f'{ip}/ppt/{f}'):
                    fld = j[0].split ('ppt/')[1]
                    if '\\' in fld:
                        fld = fld.replace ('\\', '/')

                    fl_list = natsort.natsorted(j[2])
                    for x in fl_list:
                        ext = ''.join (pathlib.Path(x).suffixes)
                        name = re.findall (r'(\w+?)(\d+)', x) [0][0]
                        count = cnt [fld] + 1
                        new_name = f'{name}{count}{ext}'
                        shutil.copy (f'{ip}/ppt/{fld}/{x}', f'{des}/{fld}/{new_name}')

                        nm[f'{f}/{x}'] = f'{fld}/{new_name}'
                        cnt[fld] = count

    # remove empty folders
    for i in os.walk (des):
        if not i[2]:
            shutil.rmtree (i[0])

    return nm, cnt


def apply_assets (assets, ft, ip, assets_cnt):
    """
    copy required assets from input to output deck 
    """
    d = OrderedDict()

    assets = natsort.natsorted (assets)
    for asset in assets:
        if '/' in asset:
            if 'slideLayouts' not in asset and 'slideMasters' not in asset and 'theme' not in asset:
                fld, f = build_names (asset)
                # if os.path.exists (f'{ft}/ppt/{fld}/{f}'):
                #     path = f'{ft}/ppt'
                #     d.update (rename (ft, fld, f, assets_cnt, path))
                # else:
                #     d.update (rename (ft, fld, f, assets_cnt))
                d.update (rename (ft, ip, fld, f, assets_cnt))
    return d


def create_json (fl, name):
    """
    creates a json file
    """
    obj = json.dumps(fl, indent=4)
    with open(f"new_json/{name}.json", "w") as outfile:
        outfile.write(obj)
    return


def apply_rels (rels, ft, assets_name):
    """
    update latest .rels files
    refactor the names of the assests and update the content
    """
    old = natsort.natsorted ([i for i in assets_name.keys()]) # old_files
    path = f'{ft}/ppt'
    
    for rel in rels:
        root, tree = build_tree (f'{path}/{rel}')
        for relation in root:
            attrib = relation.attrib
            if attrib.get ('Target')[3:] in old:
                relation.set ('Target', f"../{assets_name[attrib.get ('Target')[3:]]}")
    
        tree.write (f'{path}/{rel}', pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    
    return rels


def build_relations (pxr, ft, pp, input_deck, ss):
    """
    generate two lists of tags, one for input deck and other for output deck
    """
    root1,_ = build_tree (pxr)
    root2,_ = build_tree (f'{ft}/ppt/_rels/presentation.xml.rels')
    
    ts = total_slides (pp) # total slides
    fs_id = first_slide (pxr) # first slide id

    assest_rels = {}
    itags = []
    otags = []

    for relation in root1:
        attrib = relation.attrib
        cur_rId = int(attrib.get ('Id').split ('Id')[-1])
        if (fs_id > cur_rId) or (cur_rId > (fs_id + ts - 1)):
            itags.append(attrib["Target"])
            assest_rels[attrib['Target']] = [relation.tag, attrib['Id'], attrib['Type'], attrib['Target']]
        if not ss:
            if (fs_id <= cur_rId) and (cur_rId < (fs_id + ts)):
                itags.append(attrib["Target"])
                assest_rels[attrib['Target']] = [relation.tag, attrib['Id'], attrib['Type'], attrib['Target']]

    if ss:
        for id in ss:
            slide = f'slide{str(id)}.xml'
            for relation in root1:
                attrib = relation.attrib
                if slide in attrib['Target'] and "http" not in attrib['Target']:
                    itags.append(attrib['Target'])
                    assest_rels[attrib['Target']] = [relation.tag, attrib['Id'], attrib['Type'], attrib['Target']]

    for relation in root2:
        attrib = relation.attrib
        otags.append(attrib['Target'])
    
    itags = natsort.natsorted(itags)
    otags = natsort.natsorted(otags)

    return itags, otags, assest_rels


def update_assest_rels (nm, assest_rels):
    """
    update assest_rels by removing '../' from the target
    """
    d3_keys = [i for i in assest_rels.keys()]
    out_keys = natsort.natsorted([i for i in d3_keys])

    arels = OrderedDict(assest_rels)
    
    for i in out_keys:
        if '/' in i:
            val = arels[i]
            if '../' in i:
                val[3] = f'../{nm[i[3:]]}'
            else:
                val[3] = nm[i]
            arels[i] = val
    return arels


def update_rId (assets_cnt, itags, assest_rels):
    """
    update the rIds for assets, largest rId of the output deck
    and map old and new rIds
    """
    d1 = OrderedDict()
    cnt = OrderedDict(assets_cnt)
    max_rId = cnt['rId']
    for i in itags:
        max_rId += 1
        val = assest_rels[i]
        d1[val[1]] = f'rId{max_rId}'
        val[1] = f'rId{max_rId}'
        assest_rels[i] = val

    cnt['rId'] = max_rId
    
    return cnt, assest_rels, d1


def process_pxr (ft, assest_rels, itags):
    """
    add assests relationship in presentation.xml.rels
    """
    path = f'{ft}/ppt/_rels/presentation.xml.rels'
    root, tree = build_tree (path)
    
    for i in itags:
        val = assest_rels[i]
        tag, Id, Type, target = val
        ele = etree.Element(tag)
        etree.SubElement (root, tag, Id=Id, Type=Type, Target=target)
    
    tree.write (path, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    
    return assest_rels


def gen_xml_tag (inp_tag, out_tag):
    """
    generate a dictionary with new_tag as key and prev_tag as value
    """
    tag_dict = OrderedDict()
    
    for i in range(len(inp_tag)):
        if inp_tag[i] not in out_tag:
            tag_dict[inp_tag[i]] = [inp_tag[i-1]]
    
    return tag_dict


def create_tags (ft, tag_dict, o_tree):
    """
    create empty tags in presentation.xml
    """
    for i, o in tag_dict.items():
        subtag1 = o_tree.find (o[0])
        subtag2 = etree.Element(i)
        subtag1.addnext (subtag2)
    
    o_tree.write (f'{ft}/ppt/presentation.xml', pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)

    return 


def modify_assets_cnt (d1, d2):
    """
    modify the subtag dictionary
    """
    val_list = [i for i in d2.keys()]

    for key in val_list:
        for i in range(len(d2[key])):
            try:
                val = d1[d2[key][i][2]]
                d2[key][i][2] = val
                if None in d2[key][i]:
                    d2[key][i].remove(None)
            except:
                pass
    return d2


def add_extLst (src_xml, des_xml, ext_lst, tag_dict):
    """
    adding extlst subelements in presentation.xml file
    """
    inp_root,_ = build_tree(src_xml)
    out_root, out_tree = build_tree(des_xml)

    for relation in ext_lst:
        for elt in inp_root.findall(relation):
            out_root.append(elt)

    out_tree.write(des_xml, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    return


def get_prep_tags (src_xml, d1):
    """
    create a dict with tags as keys and subtags as values
    """
    root, tree = build_tree(src_xml)
    d2 = OrderedDict()
    
    # list of old rIds
    rId_lis = [i for i in d1.keys()]
    
    nmsps =  root.nsmap['r']
    ext_lst = []

    for relation in root:
        for ele in relation:
            attrib = ele.attrib
            tag = ele.tag
            
            if attrib.get(f"{{{nmsps}}}id"):
                if attrib.get(f"{{{nmsps}}}id") in rId_lis:
                    if relation.tag in d2:
                        val = d2[relation.tag]
                        val.append([tag, attrib.get('id'), attrib.get(f"{{{nmsps}}}id")])
                        d2[relation.tag] = val
                    else:
                        d2[relation.tag] = [[tag, attrib.get('id'), attrib.get(f"{{{nmsps}}}id")]]
            else:
                if 'uri' in ele.attrib:
                    if relation.tag not in ext_lst:
                        ext_lst.append(relation.tag)
    
    d2 = modify_assets_cnt (d1, d2)
    
    return d2, ext_lst


def add_subtags (ft, path, pxml_subtags):
    """
    add subtags in the presentation.xml file
    """
    root, tree = build_tree (path)
    nmsps =  root.nsmap['r']
    for k,v in pxml_subtags.items():
        subtag1 = tree.find (k)
        
        for i in v:
            if 'rId' not in i[1]:
                rId = f"{{{nmsps}}}id"
                subtext = etree.SubElement(subtag1, i[0])
                subtext.attrib['id'] = i[1]
                subtext.attrib[rId] = i[2]
            else:
                subtext = etree.SubElement(subtag1, i[0])
                subtext.attrib[rId] = i[1]
    
    tree.write(path, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    return


def clean_prep_xml (des_xml, rels_rIds, pxml_subtags):
    """
    remove extra subtags if their rId is not present in 
    the presentation.xml.rels file
    """
    root, tree = build_tree(des_xml)
    nmsps =  root.nsmap['r']
    rId = f"{{{nmsps}}}id"
    for k,v in pxml_subtags.items():
        subtag1 = tree.find(k)
        for i in subtag1:
            if i.attrib.get(rId):
                if i.attrib.get(rId) not in rels_rIds:
                    subtag1.remove(i)

    tree.write(des_xml, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    return


def process_dup_rel (ft):
    """
    remove the duplicates entries of 'Target' field 
    form presentation.xml.rels if any and 
    generate a dict of key:target, val:rId 
    """
    path = f'{ft}/ppt/_rels/presentation.xml.rels'
    root, tree = build_tree(path)
    
    rd = OrderedDict() # relation dictionary
    for relation in root:
        rIds = []
        attrib = relation.attrib
        if attrib['Target'] in rd.keys():
            val = rd[attrib['Target']]
            val.append(attrib['Id'])
            rd[attrib['Target']] = val
        else:
            rd[attrib['Target']] = [attrib['Id']]

    # getting duplicates rIds
    dup_rIds = []
    for k,v in rd.items():
        if len(v) > 1:
            dup_rIds.append(v.pop(0))
            rd[k] = v

    # removing relation
    for relation in root:
        attrib = relation.attrib
        if attrib['Id'] in dup_rIds:
            root.remove(relation)

    r_rIds = [relation.attrib['Id'] for relation in root]

    tree.write(path, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    
    return rd, r_rIds


def scan_sldsz (src_xml, des_xml):
    """
    scan slide size for templating in presentation.xml
    """
    _,i_tree = build_tree (src_xml)
    _,o_tree = build_tree (des_xml)

    tag = "{http://schemas.openxmlformats.org/presentationml/2006/main}sldSz"

    inp_sldsz = i_tree.find (tag).attrib
    cx = inp_sldsz['cx']
    cy = inp_sldsz['cy']
    Type = inp_sldsz.get('type')
    out_sldsz = o_tree.find (tag).attrib
    
    if int(out_sldsz['cx']) < int(cx):
        if int(out_sldsz['cy']) < int(cy):
            out_sldsz['cx'] = cx
            out_sldsz['cy'] = cy
        else:
            out_sldsz['cx'] = cx
    
    if not Type and out_sldsz.get('type'):
        del out_sldsz['type']
    
    o_tree.write(des_xml, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    
    return


def process_pr (ft, ip, rd, rels_rIds):
    """
    update the presentation.xml file
    """
    src_xml = f'{ip}/ppt/presentation.xml'
    des_xml = f'{ft}/ppt/presentation.xml'

    inp_root, inp_tree = build_tree (src_xml)
    out_root, out_tree = build_tree (des_xml)

    inp_tag = [relation.tag for relation in inp_root]
    out_tag = [relation.tag for relation in out_root]

    tag_dict = gen_xml_tag (inp_tag, out_tag)
    create_tags (ft, tag_dict, out_tree)

    pxml_subtags, ext_lst = get_prep_tags (src_xml, rd)
    # add_extLst(src_xml, des_xml, ext_lst, tag_dict)
    create_json (pxml_subtags, '05_pxml_subtag')
    clean_prep_xml (des_xml, rels_rIds, pxml_subtags)
    add_subtags (ft, des_xml, pxml_subtags)
    scan_sldsz (src_xml, des_xml)
    
    return


def pres_configs (pxr, ft, ip, pp, input_deck, ss, assets_name, assets_cnt):
    """
    deals with rels and xml file of presentation
    """
    nm = dict (assets_name)
    cnt = dict (assets_cnt)
    
    itags, otags, assest_rels = build_relations (pxr, ft, pp, input_deck, ss)
    itags = remove_dup (itags, assest_rels, otags)
    
    rel = update_assest_rels (nm, assest_rels)
    assest_rels.update (rel)
    
    cnt, rel, rd = update_rId(cnt, itags, assest_rels)
    assest_rels.update (rel)
    
    process_pxr (ft, assest_rels, itags) # process presentation.xml.rels
    
    prep_rels_rIds, rels_rIds = process_dup_rel (ft)
    
    create_json (assest_rels, '03_prepRelSubtag')
    create_json (prep_rels_rIds, '04_prep_rels_rIds')

    process_pr (ft, ip, rd, rels_rIds) # process presentation.xml
    return cnt


def handle_configs (ip, ft):
    """
    handle configuration files
    """
    inp_path = '/'.join([ip, 'ppt'])
    out_path = f'{ft}/ppt'

    config_fls = [i for i in os.listdir(inp_path) if os.path.isfile(f'{inp_path}/{i}')]
    
    mergables = ['commentAuthors.xml', 'tableStyles.xml']
    sing_prop = ['viewProps.xml', 'presProps.xml']
    ignore = ['revisionInfo.xml']

    for i in config_fls:
        inp_fl = f'{inp_path}/{i}'
        out_fl = f'{out_path}/{i}'
        
        if os.path.isfile(f'{out_path}/{i}'):
            root1,tree1 = build_tree(inp_fl)
            root2,tree2 = build_tree(out_fl)
            if i in mergables:
                try:
                    for relation in [f"{root1[0].tag}"]:
                        for elt in root1.findall(relation):
                            root2.append(elt)
                except:
                    pass
            elif i in sing_prop:
                if i == 'presProps.xml':
                    inp_d = {}
                    out_lis = []
                    nm = root1.nsmap['p']
                    tag0 = f"{{{nm}}}extLst"
                    for relation in [f"{root1[0].tag}"]:
                        fp = root1.find(tag0)
                        for ele in fp:
                            attrib = ele.attrib
                            if attrib['uri'] not in inp_d.keys():
                                inp_d[attrib['uri']] = ele

                    for relation in [f"{root2[0].tag}"]:
                        fp = root2.find(tag0)
                        for ele in fp:
                            attrib = ele.attrib
                            out_lis.append(attrib['uri'])
                    for k,v in inp_d.items():
                        if k not in out_lis:
                            tag1 = root2.find(tag0)
                            tag1.append(v)
            else:
                pass
            tree2.write(out_fl, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
        else:
            shutil.copyfile(inp_fl, out_fl)


    return


def handle_cleaning (ft):
    """
    remove unnecessary files like changesInfos, printerSettings
    """
    extra_fl = ['changesInfos', 'printerSettings']
    fld_path = f'{ft}/ppt'
    out_pxr = f'{fld_path}/_rels/presentation.xml.rels'
    root, tree = build_tree(out_pxr)

    for i in extra_fl:
        path = f'{fld_path}/{i}'
        if os.path.isfile(path):
            shutil.rmtree(path)

        for relation in root:
            attrib = relation.attrib
            if i in attrib['Target']:
                root.remove(relation)

    tree.write(out_pxr, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    
    return out_pxr


def content(ft, ip, assets_name):
    """
    add content_type of new assets in [Contant_Types].xml file
    """

    f = '[Content_Types].xml'
    inp_path = '/'.join([ip, f])
    out_path = '/'.join([ft, f])

    cnt_lst = []
    asset_lst = []
    def_att = []
    d = dict()

    root1,tree1 = build_tree(inp_path)
    root2,tree2 = build_tree(out_path)

    # get all the extensions belongs to "Default" tag
    for relation in root2:
        if 'Default' in relation.tag:
            def_att.append(relation.attrib['Extension'])
        else:
            break

    for relation in root1:
        if 'Override' in relation.tag:
            attrib = relation.attrib['PartName'][1:]
            try:
                cnt = attrib.split('ppt/')[-1]
                ini = '/ppt/'
            except:
                cnt = attrib
                ini = '/'
            
            if cnt in assets_name.keys():
                relation.attrib['PartName'] = f'{ini}{assets_name[cnt]}'
                cnt_lst.append(relation)
                # asset_lst.append(relation.attrib['PartName'])
            else:
                cnt_lst.append(relation)
            
            if relation.attrib['PartName'] not in asset_lst:
                asset_lst.append(relation.attrib['PartName'])
        
        else:
            attrib = relation.attrib['Extension']
            if attrib not in def_att:
                cnt_lst.append(relation)

    cnt_lst = natsort.natsorted(cnt_lst)
    for ele in cnt_lst:
        prev = tree2.find(ele.tag)
        prev.addnext(ele)

    tree2.write(out_path, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)

    unq_attr = []
    for relation in root2:
        if 'Override' in relation.tag:
            if relation.attrib['PartName'] not in unq_attr:
                unq_attr.append(relation.attrib['PartName'])
            else:
                root2.remove(relation)
    
    tree2.write(out_path, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)

    return out_path


def deck_handler(r, msg, bp, d, assets_cnt):
    """
    handle the deck creation process
    """
    input_deck = msg ['d']
    
    ip = f'{bp}/tmp/{str(r)}/{input_deck}'
    op = f'{bp}/output/{str(r)}'
    pp = f'{bp}/presentations/{input_deck}.pptx'
    
    try:
        os.makedirs (op)
        os.makedirs (ip)
    except:
        pass
    
    a = unpack (pp, ip)
    ft = op
    
    make_dir (ip, ft, input_deck)
    
    if d == 1:
        make_structure (ip, ft, input_deck)

    ss = msg ['s']
    if not ss:
        ss = build_slides (pp)    
    
    assets = []
    pxr = f'{ip}/ppt/_rels/presentation.xml.rels'
    assets = build_assets (pxr, input_deck, assets, ss)
    
    assets_name = apply_assets (assets, ft, ip, assets_cnt)
    
    nm, cnt = apply_mend_assets (ip, ft, d, assets_name, assets_cnt)
    
    assets_name.update (nm)
    assets_cnt.update (cnt)

    rels = build_rels (assets_name)
    apply_rels (rels, ft, assets_name)
    
    # updating rels and xml files of presentation
    cnt = pres_configs (pxr, ft, ip, pp, input_deck, ss, assets_name, assets_cnt)
    assets_cnt.update (cnt)
    
    # adding contents in the Content-Type file
    content (ft, ip, assets_name)
    
    # handling the properties files
    handle_configs (ip, ft)
    
    # removing extra files
    handle_cleaning (ft)

    # creating json files
    create_json(assets_name, '01_refactored_names')
    create_json(assets_cnt, '02_refactoring_count')

    return assets_cnt

    
def deck_render_effect ():
    # msg = [41, {'d': 'Onboarding', 's':  [2,4,6]}]
    msg = [41, {'d': 'Onboarding', 's':  [2,4,6]}, {'d': 'Presentation1','s':  None}]
    r = msg.pop(0)
    
    bp = os.getcwd ()
    emp = f'{bp}/output/{str(r)}'
    new(emp)
    assets_cnt = OrderedDict(max_rId(emp))
    
    d = 1
    while msg:
        assets_cnt = deck_handler(r, msg.pop(0), bp, d, assets_cnt)
        d += 1
    
    pack (emp, "Test")

deck_render_effect()
