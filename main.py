import os
import shutil
import json
import xmltodict
import natsort
import zipfile
import re
import pathlib

from pptx import Presentation
from lxml import etree
from collections import OrderedDict


def unzip(src, des):
    """
    unzip the deck
    """
    with zipfile.ZipFile(src, 'r') as zip_ref:
        zip_ref.extractall(des)
    return


def zipdir(path, file_name):
    """
    zip extracted deck to get output deck
    """
    length = len(path)
    zipf = zipfile.ZipFile('output/'+f'Test_{file_name}.pptx', 'w', zipfile.ZIP_DEFLATED)
    for root, dirs, files in os.walk(path):
        folder = root[length:] # path without "parent"
        for file in files:
            zipf.write(os.path.join(root, file), os.path.join(folder, file))
    zipf.close()
    return


def gen_tree(path):
    """
    pass the path of the xml document to enable the parsing process
    """
    print("CALLING.. Tree")
    parser = etree.XMLParser(remove_blank_text=True)
    tree = etree.parse(path, parser)
    root = tree.getroot()    
    return root, tree


def max_rId():
    """
    returns maximum rId
    """
    path = f'{output_path}/ppt/_rels/presentation.xml.rels'
    root, tree = gen_tree(path)
    
    rIds = []
    
    for relation in root:
        attrib = relation.attrib
        rId = int(attrib.get('Id').split('Id')[-1])
        rIds.append(rId)
    return {'rId': max(rIds)}


def xml_to_dict(path):
    """
    convert xml to dict
    """
    with open(path) as xml_file:
        data_dict = xmltodict.parse(xml_file.read())
        xml_file.close()
    if isinstance(data_dict["Relationships"]["Relationship"], list):
        data = sorted(data_dict["Relationships"]["Relationship"], key=lambda item: int(item['@Id'].split('Id')[1]))
    else:
        data = [data_dict["Relationships"]["Relationship"]]
    return data


def total_slides(path):
    """
    returns total number of slides
    """
    print("CALLING.. total_slides")
    prs = Presentation(path)
    tot_slides = len(prs.slides._sldIdLst)
    return tot_slides


def first_slide(path):
    """
    returns the first slide rId
    """
    print("CALLING.. First_slide")
    root, _ = gen_tree(path)

    slide = 'slide1.xml'
    for relation in root:
        attrib = relation.attrib
        if slide in attrib['Target']:
            return int(attrib['Id'].split('Id')[-1])


def get_rels(dict_1):
    """
    list latest .rels files
    """
    # a = dict_1.keys()
    a = [i for i in dict_1.values()]
    lis = natsort.natsorted([i for i in a if '_rels' in i])
    return lis


def remove_dup(files1, dict_3, files2):
    """
    remove duplicates from the rels files
    """
    l1 = files1[:]
    for i in l1:
        if '/' not in i:
            if i in files2:
                files1.remove(i)
                del dict_3[i]
    return files1   


def new(path):
    """
    create, move and unzip the empty output deck
    """
    fq_empty = "resources/Empty.pptx"
    # create
    prs = Presentation()
    prs.save(fq_empty)
    # move
    d = ".".join ([path, "pptx"]) # "output/41.pptx"
    shutil.move (fq_empty, d)
    # unzip
    unzip (d, d.split('.')[0])
    m_rId = max_rId()
    return m_rId


def make_dir(file_name): # output_file_loc = des
    """
    creates input deck directories which does not exists in the output deck
    """
    for i in os.walk(f'{tmp_path}/{file_name}'):
        fld = i[0].split(file_name)[-1]
        if fld:
            loc = f"{output_path}{fld}"
            if not os.path.exists(f'{output_path}/{fld}'):
                os.makedirs(f'{output_path}/{fld}')
    print("MAKE_DIR completed...")        
    return


def make_structure(file_name):
    """
    creates structure of input deck
    """
    for i in os.walk(f'{tmp_path}/{file_name}'):
        fld = i[0].split(file_name)[-1]
        if fld:
            loc = f"{output_path}{fld}"
            if 'ppt' not in loc and (file_name not in loc):
                shutil.rmtree(f'{output_path}/{fld}')
                shutil.copytree(f'{tmp_path}/{file_name}/{i[0].split(file_name)[-1]}', f'{output_path}/{fld}')
            # elif 'ppt' in loc:
            #     if 'slideLayouts' in loc or 'slideMasters' in loc or 'theme' in loc:
            #         # copy_mandatory()
            #         pass
            
            # # if not os.path.exists(f'{output_file_loc}/{fld}'):
            # #     if 'ppt' not in f'{output_file_loc}/{fld}':
            # #         shutil.copytree(f"{tmp_path}/{file}/{i[0].split(f'{file}/')[-1]}", f'{output_file_loc}/{i[0].split(file)[-1]}')
            # #     else:
            # #         os.makedirs(f'{output_file_loc}/{i[0].split(file)[-1]}')
    return


def add_files(path, file_name, target_files, slides=None):
    """
    returns a list of files that needs to be modified in output deck
    """
    print("CALLING.. Add_files")
    data = xml_to_dict(path)
    if slides:
        global sldIds
        
        tot_slides = total_slides(f'{input_decks}/{file_name}.pptx')
        first_slide_id = first_slide(path)
        
        files = []
        for i in data:
            current_rId = int(i['@Id'].split('Id')[1])
            if (first_slide_id > current_rId) or (current_rId > (first_slide_id+tot_slides-1)):
                if 'slideLayouts' not in i['@Target'] and 'slideMasters' not in i['@Target'] and 'theme' not in i['@Target']:
                    files.append(i['@Target'])
        for i in files:
            if '/' in i:
                a = i.split('/')
                fld, fl = a[0], a[1]
                if os.path.exists(f'{tmp_path}/{file_name}/ppt/{fld}/_rels'):
                    if os.path.isfile(f'{tmp_path}/{file_name}/ppt/{fld}/_rels/{fl}.rels'):
                        target_files.append(f'{fld}/_rels/{fl}.rels')
                    # fl_list = os.listdir(f'{tmp_path}/{file_name}/ppt/{fld}/_rels')
                    # for i in fl_list:
                    #     if 
                    
            pass# if rel exists then will add it to list
        
        target_files = target_files + files
        
        for id in slides:
            slide = f'slide{str(id)}.xml'
            sldIds.append([i["@Id"] for i in data if slide in i["@Target"] and "http" not in i["@Target"]][0])
            target_files.append([i["@Target"] for i in data if slide in i["@Target"] and "http" not in i["@Target"]][0])
            target_files.append(f'slides/_rels/{slide}.rels')
            add_files(f'{tmp_path}/{file_name}/ppt/slides/_rels/{slide}.rels', file_name, target_files)
    else:
        for i in data:
            # handling duplicacy
            if i["@Target"] in target_files or i["@Target"][3:] in target_files:
                pass
            elif "http" not in i["@Target"]:
                # if 'slideLayouts' not in i['@Target'] and 'slideMasters' not in i['@Target'] and 'theme' not in i['@Target']:
                if '../' in i["@Target"]:
                    target_files.append(i['@Target'][3:])
                else:
                    target_files.append(i['@Target'])
                
                if ".." in i['@Target'] and "xml" in i['@Target']:
                    path = f"{tmp_path}/{file_name}/ppt/{i['@Target'].split('..')[1].split('/')[1]}/_rels/{i['@Target'].split('..')[1].split('/')[2]}.rels"
                    
                    if os.path.exists(path):
                        # handling rels files
                        target_files.append(path.split('ppt/')[1])
                        add_files(path, file_name, target_files)
                        
        # print("UUU2: ", target_files)
    # print("UUU: ", target_files)
    return target_files


def get_fld_fl(file):
    """
    returns folder anme and file name
    """
    if '_rels' in file: # slides/_rels/slide2.xml.rels
        sp = file.split('/')
        fl_name = sp[-1]
        fld_name = f'{sp[0]}/{sp[1]}'
    elif '../' in file:
        _,fld_name,fl_name = file.split('/')
    else:
        fld_name,fl_name = file.split('/')
    
    return fld_name, fl_name


def list_target(target_files, d2):
    """
    creates a dict with number of files
    """
    # d2 = OrderedDict()
    count = 0
    for file in target_files:
        if '/' in file:
        # get folder and file name
            fld, fl = get_fld_fl(file)
            if fld not in d2:
                d2[fld] = 0
        if 'slideMasters' not in d2:
            d2['slideMasters'] = 0
        if 'slideLayouts' not in d2:
            d2['slideLayouts'] = 0
        if 'theme' not in d2:
            d2['theme'] = 0
    return d2
    

def rename(path, fld, fl, dict_2): # fld=media, fl=image1.png
    """
    rename a file
    """
    d1 = OrderedDict()
    
    ext = ''.join(pathlib.Path(fl).suffixes)
    name = re.findall(r'(\w+?)(\d+)', fl)[0][0]
    
    count = dict_2[fld]+1
    new_name = f'{name}{count}{ext}'
    # if 'slideMasters' in fld:
        # print("NEW_NAME", new_name)
        # print("GGG: ", dict_2)
    if 'ppt' in path:
        shutil.copy(f'{path}/{fld}/{fl}', f"{output_path}/ppt/{fld}/{new_name}")
    else:
        shutil.copy(f'{path}/{fld}/{fl}', f"{output_path}/{fld}/{new_name}")
    d1[f'{fld}/{fl}'] = f'{fld}/{new_name}'
    dict_2[fld] = count
    # print("RENAME: ", fld, fl, new_name)
    return d1


def del_files(rels_fl, last_fl, path):
    """
    delete extra files
    """
    for i in rels_fl:
        if i[:-5] not in last_fl:
            os.remove(f'{path}/{i}')


def copy_mandatory(src, des, deck, dict_1):
    """
    copy mandatory files
    """
    # print("SRC: ", src, "\nDES: ", des)
    print("COPY MANDATORY CALLING")
    m_list = ['slideLayouts', 'theme', 'slideMasters']
    # d1 = OrderedDict()
    if deck == 1:
        for fl in m_list:
            count = 0
            for i in os.walk(f'{src}/ppt/{fl}'):
                count = len(i[2])
            if os.path.exists(f'{des}/{fl}'):
                shutil.rmtree(f'{des}/{fl}')
            shutil.copytree(f'{src}/ppt/{fl}', f'{des}/{fl}')
            dict_2[fl] = count
            if os.path.exists(f'{src}/ppt/{fl}/_rels'):
                dict_2[f'{fl}/_rels'] = count
            
            for i in os.walk(f'{src}/ppt/{fl}'):
                fld = i[0].split('ppt/')[1]
                fl_list = natsort.natsorted(i[2])
                for j in fl_list:
                    dict_1[f'{fld}/{j}'] = f'{fld}/{j}'
    else:
        # print("SSSS: ", dict_2)
        for i in m_list:
            if os.path.exists(f'{des}/{i}'):
                for j in os.walk(f'{src}/ppt/{i}'):
                    # j = (path, [folder], [file])
                    fld = j[0].split('ppt/')[1]
                    
                    # print("TTT: ", fld, "--", j[2]) # array
                    fl_list = natsort.natsorted(j[2])
                    for x in fl_list:
                        ext = ''.join(pathlib.Path(x).suffixes)
                        name = re.findall(r'(\w+?)(\d+)', x)[0][0]
                        count = dict_2[fld]+1
                        new_name = f'{name}{count}{ext}'
                        # print("UUUU: ", new_name, count)
                        # print("UUUU: ", f'{src}/ppt/{i}/{x}')
                        # print("1111: ", f'{des}/{fld}/{new_name}')
                        shutil.copy(f'{src}/ppt/{fld}/{x}', f'{des}/{fld}/{new_name}')
                        
                        dict_1[f'{i}/{x}'] = f'{fld}/{new_name}'
                        dict_2[fld] = count
    
    # remove empty folders
    for i in os.walk(des):
        if not i[2]:
            shutil.rmtree(i[0])

    
    print("MANDATORY DONE!!")
    
    
    return dict_1, dict_2


def copy_target(target_files, file_name, tmp_loc, dict_2):
    """
    copy target files from tmp to output folder 
    """
    d_1 = OrderedDict()
    
    target_files = natsort.natsorted(target_files)
    for i in target_files:
        if '/' in i:
            if 'slideLayouts' not in i and 'slideMasters' not in i and 'theme' not in i:
                fld_name,fl_name = get_fld_fl(i)
                # print("DDDD: ", fl_name)
                if os.path.exists(f'{tmp_loc}/ppt/{fld_name}/{fl_name}'):
                    path = f'{tmp_loc}/ppt'
                    d_1.update(rename(path, fld_name, fl_name, dict_2))
                else:
                    d_1.update(rename(tmp_loc, fld_name, fl_name, dict_2))
        else:
            if not os.path.isfile(f'{output_path}/ppt/{i}'):
                shutil.copyfile(f'{tmp_loc}/ppt/{i}', f'{output_path}/ppt/{i}')
    
    return d_1


def create_json(fl, name):
    """
    creates a json files
    """
    obj = json.dumps(fl)
    with open(f"new_json/{name}.json", "w") as outfile:
        outfile.write(obj)
    return


def update_rels(fl_list, tmp_loc, dict_1):
    """
    update latest .rels files
    """
    old_files = natsort.natsorted([i for i in dict_1.keys()])
    path = f'{output_path}/ppt'
    for i in fl_list:
        root, tree = gen_tree(f'{path}/{i}')
        for relation in root:
            attrib = relation.attrib
            if attrib.get('Target')[3:] in old_files:
                relation.set('Target', dict_1[attrib.get('Target')[3:]])
        tree.write(f'{path}/{i}', pretty_print=True, xml_declaration=True, encoding='UTF-8')
    return


def get_relations(inp_path, file_name, slides):
    """
    returns f1:list of inp targets, d3:dict of targets, 
    f2: list of out target and sldIds: list of rIds of slides
    """
    root1,_ = gen_tree(inp_path)
    root2,_ = gen_tree(f'{output_path}/ppt/_rels/presentation.xml.rels')
    data = xml_to_dict(inp_path)
    tot_slides = total_slides(f'{input_decks}/{file_name}.pptx')
    first_slide_id = first_slide(inp_path)

    dict_3 = OrderedDict()
    files1 = []
    files2 = []
    sldIds = []

    for relation in root1:
        attrib = relation.attrib
        # print("ATTRIB: ", attrib)
        current_rId = int(attrib.get('Id').split('Id')[-1])
        if (first_slide_id > current_rId) or (current_rId > (first_slide_id+tot_slides-1)):
            files1.append(attrib["Target"])
            dict_3[attrib['Target']] = [relation.tag, attrib['Id'], attrib['Type'], attrib['Target']]
        if not slides:
            if (first_slide_id <= current_rId) or (current_rId < (first_slide_id+tot_slides+1)):
                files1.append(attrib["Target"])
                sldIds.append(attrib['Id'])
                dict_3[attrib['Target']] = [relation.tag, attrib['Id'], attrib['Type'], attrib['Target']]

    if slides:
        for id in slides:
            slide = f'slide{str(id)}.xml'
            for relation in root1:
                attrib = relation.attrib
                # print("TYPE: ", type(attrib['Type']))
                if slide in attrib['Target'] and "http" not in attrib['Target']:
                    files1.append(attrib['Target'])
                    sldIds.append(attrib['Id'])
                    dict_3[attrib['Target']] = [relation.tag, attrib['Id'], attrib['Type'], attrib['Target']]
    files1 = natsort.natsorted(files1)

    for relation in root2:
        attrib = relation.attrib
        files2.append(attrib['Target'])
    files2 = natsort.natsorted(files2)
    
    return files1, dict_3, files2, sldIds


def update_dict_3(dict_1, dict_3):
    """
    update dict_3
    """
    inp_keys = [i for i in dict_1.keys()]
    d3_keys = [i for i in dict_3.keys()]
    out_keys = natsort.natsorted([i for i in d3_keys])
    # print("UUUU: ", out_keys)
    
    for i in out_keys:
        if '/' in i:
            val = dict_3[i]
            if '../' in i:
                val[3] = f'../{dict_1[i[3:]]}'
            else:
                val[3] = dict_1[i]
            dict_3[i] = val
    return dict_3


def update_rId(dict_2, files1, dict_3):
    """
    update the rIds
    """
    max_rId = dict_2['rId']
    for i in files1:
        max_rId += 1
        val = dict_3[i]
        val[1] = f'rId{max_rId}'
        dict_3[i] = val
        
    dict_2['rId'] = max_rId
    return dict_2, dict_3


def write_rels(dict_3, files1):
    """
    adding assests in presentation.xml.rels
    """
    path = f'{output_path}/ppt/_rels/presentation.xml.rels'
    root, tree = gen_tree(path)
    for i in files1:
        val = dict_3[i]
        tag, Id, Type, target = val
        ele = etree.Element(tag)
        etree.SubElement(root, tag, Id=Id, Type=Type, Target=target)
    tree.write(path, pretty_print=True, xml_declaration=True, encoding='UTF-8')
        

def xml_tag(inp_tag, out_tag):
    """
    returns a dict, key=new_tag, value=prev_tag
    """
    tag_dict = OrderedDict()
    sub_tag = OrderedDict()
    for i in range(len(inp_tag)):
        if inp_tag[i] not in out_tag:
            # if i == 0:
            #     tag_dict[inp_tag[i]] = [0]
            # else:
            tag_dict[inp_tag[i]] = [inp_tag[i-1]]
    return tag_dict


def create_tags(inp_root, out_root, tag_dict, i_tree, o_tree):
    """
    modify content of presentation.xml
    """
    for i, o in tag_dict.items():
        # if o[0] == 0:
        #     pass
        subtag1 = o_tree.find(o[0])
        subtag2 = etree.Element(i)
        # for ele in elements:
        #     subtext = etree.SubElement(subtag2, ele)
        # subtext = etree.SubElement(subtag2)
        subtag1.addnext(subtag2)
            
    with open (f'{output_path}/ppt/presentation.xml', 'wb') as f:
        f.write(etree.tostring(out_root, pretty_print = True))
    
    return


def sub_tag(tag_dict, inp_root, inp_tree):
    """
    lists subtags of missing tags
    """
    dict_4 = OrderedDict()
    
    tag = [i for i in tag_dict.keys()]
    # print("TAG: ", tag)
    for i in tag:
        a = inp_tree.find(i)
        attrib = a.attrib
        # print("VAL: ", attrib)
        # print("TYPE: ", type(attrib))
        # fi=or 
        # print("VAL: ", eval(str(attrib))
        # print("VAL: ", attrib.keys())
        # for k in attrib:
        #     print("VAL: ", k)
        # if 
        for j in a:
            # dict_4[attrib['']]
            # dict_3[attrib['Target']] = [relation.tag, attrib['Id'], attrib['Type'], attrib['Target']]
            # print("JJ: ", j.attrib)
            pass


def filter_sld_rIds(d1, sldIds):
    """
    filter slides on required slide rIds
    """
    key = "{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst"
    a = d1[key]
    b = [i for i in a if i[-1] in sldIds]
    d1[key] = b
    return d1


def filter_emp_tags(d1, tag_dict):
    """
    filter empty tags, tags with no Id and rId
    """
    for i in tag_dict:
        pass
    # working
      

def get_prep_tags(src_xml, tag_dict, sldIds):
    """
    get tags and subtags form input presentation.xml
    """
    root, tree = gen_tree(src_xml)
    
    d1 = OrderedDict()
    tag_list = [i for i in tag_dict.keys()]
    nmsps =  root.nsmap['r']
    
    for relation in root:
        if relation.tag in tag_list:
            for ele in relation:
                attrib = ele.attrib
                tag = ele.tag
                if relation.tag in d1:
                    # print("IF : ", relation.tag)
                    try:
                        val = d1[relation.tag]
                        val.append([tag, attrib.get('id'), attrib.get(f"{{{nmsps}}}id")])
                        d1[relation.tag] = val
                    except:
                        pass
                else:
                    # print("ELSE: ", relation.tag)
                    d1[relation.tag] = [[tag, attrib.get('id'), attrib.get(f"{{{nmsps}}}id")]]
    # print("DDDD: ", d1)
    d1 = filter_sld_rIds(d1, sldIds)
    d1 = filter_emp_tags(d1, tag_dict)
    create_json(d1, 'dict_4')
    return

def write_pres(tmp_loc, sldIds):
    """
    update the presentation.xml file
    """
    src_xml = f'{tmp_loc}/ppt/presentation.xml'
    des_xml = f'{output_path}/ppt/presentation.xml'
    
    inp_root, inp_tree = gen_tree(src_xml)
    out_root, out_tree = gen_tree(des_xml)
    
    inp_tag = [relation.tag for relation in inp_root]
    out_tag = [relation.tag for relation in out_root]
    
    # print("ITAG: ", inp_tag, "\nOUT_TAG: ", out_tag)
    tag_dict = xml_tag(inp_tag, out_tag)
    
    create_tags(inp_root, out_root, tag_dict, inp_tree, out_tree)
    get_prep_tags(src_xml, tag_dict, sldIds)
    sub_tag(tag_dict, inp_root, inp_tree)


def rel_duplicates():
    """
    remove the duplicates entries if any
    """
    path = f'{output_path}/ppt/_rels/presentation.xml.rels'
    root, tree = gen_tree(path)
    d1 = OrderedDict()
    for relation in root:
        rIds = []
        attrib = relation.attrib
        if attrib['Target'] in d1.keys():
            val = d1[attrib['Target']]
            val.append(attrib['Id'])
            d1[attrib['Target']] = val
        else:
            d1[attrib['Target']] = [attrib['Id']]
    # getting duplicates rIds
    dup_rids = [v[0] for k,v in d1.items() if len(v)!=1]
    
    # removing relation
    for relation in root:
        attrib = relation.attrib
        if attrib['Id'] in dup_rids:
            root.remove(relation)
    tree.write(path, pretty_print=True, xml_declaration=True, encoding='UTF-8')
    return


def presenation_files(inp_pres_rels, file_name, slides, dict_1, dict_2, tmp_loc):
    """
    deals with rels and xml file of presentation
    """
    files1, dict_3, files2, sldIds = get_relations(inp_pres_rels, file_name, slides)
    print("SLDIDS: ", sldIds)
    files1 = remove_dup(files1, dict_3, files2)
    dict_3 = update_dict_3(dict_1, dict_3)
    dict_2, dict_3 = update_rId(dict_2, files1, dict_3)
    create_json(dict_3, 'dict_3')
    write_rels(dict_3, files1)
    rel_duplicates()
    write_pres(tmp_loc, sldIds)
                

def deck_handler(id, msg, deck, dict_2):
    """
    handle the deck and select files for output deck
    """
    file_name, slides = msg['d'], msg['s']
    target_files = []
    
    m_rId = new(output_path)
    dict_2.update(m_rId)
    tmp_loc = f'{tmp_path}/{file_name}'
    
    # unzip the input deck
    unzip(f'{input_decks}/{file_name}.pptx', tmp_loc)
    
    # creates folder structure of the input deck
    make_dir(file_name)
    inp_pres_rels = f'{tmp_loc}/ppt/_rels/presentation.xml.rels'
    
    if deck == 1:
        make_structure(file_name)
    
    target_files = add_files(inp_pres_rels, file_name, target_files, slides)
    # print("TARGET: ", target_files)
    dict_2.update(list_target(target_files, dict_2))
    dict_1 = copy_target(target_files, file_name, tmp_loc, dict_2)
    d1, d2 = copy_mandatory(tmp_loc, f'{output_path}/ppt', deck, dict_1)
    dict_1.update(d1)
    dict_2.update(d2)
    
    create_json(dict_1, 'dict_1')
    create_json(dict_2, 'dict_2')

    # modify the rels files
    rels_list = get_rels(dict_1)
    update_rels(rels_list, tmp_loc, dict_1)
    presenation_files(inp_pres_rels, file_name, slides, dict_1, dict_2, tmp_loc)


if __name__ == '__main__':
    
    base_path = os.path.dirname(os.path.realpath(__file__))
    print("CURRENT_DIR:", base_path)
    sldIds = []
    dict_1 = OrderedDict()
    dict_2 = OrderedDict()
    
    # sample_msg = [41, {'d': 'Onboarding','s':  [2,4,6]}, {'d': 'Presentation1','s':  [1]}]
    sample_msg = [41, {'d': 'Onboarding','s':  [2, 4, 6]}]
    # sample_msg = [41, {'d': 'Presentation1','s':  [1]}]
    # sample_msg = [41, {'d': 'BI Case Studies','s':  [2, 3]}]

    render_id = sample_msg.pop(0)
    
    output_path = f'{base_path}/output/{str(render_id)}'
    tmp_path = f'{base_path}/tmp/{render_id}'
    input_decks = f'{base_path}/presentations'
    
    # print("TMP_PATH:", tmp_path, '\nOUT_PATH: ', output_path)  

    try:
        os.makedirs(output_path)
        os.makedirs(tmp_path)
    except:
        print("DIR ALREADY EXIST")
    
    # iterating all the messages
    deck = 1
    while sample_msg:
        deck_handler(render_id, sample_msg.pop(0), deck, dict_2)
        deck += 1

    # zipdir(f'{output_file_loc}', "Test")
