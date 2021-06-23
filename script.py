import zipfile
import os
import shutil
import xmltodict
import json

from pptx import Presentation
from lxml import etree, objectify
from zip_unzip import unzip, zipdir
from copy_files import copy_rel, xml_to_dict, copy_mandatory


# zip_unzip.py


def ig_d(dir, files):
    """
    filter ppt folder
    """
    return [f for f in files if f=='ppt']


def ig_f(dir, files):
    """
    filter all files
    """
    return [f for f in files if os.path.isfile(os.path.join(dir, f))]
    

# copy all rel files
# def copy_rel(src, des):


# copy mandatory files
# def copy_mandatory(src, des):
    

# convert xml to dict
# def xml_to_dict(path):


def add_files(path, file_name, slides=[]):
    """
    add files from input deck to output deck
    """
    global target
    data = xml_to_dict(path)

    if slides:
        # get total slides
        prs = Presentation(dir_path+'/presentations/'+file_name+'.pptx')
        global tot_slides, first_slide_id
        tot_slides = len(prs.slides._sldIdLst)
        # get rId of first slide
        first_slide = "slide1.xml"
        first_slide_id = int([i["@Id"] for i in data if first_slide in i['@Target']][0].split('Id')[1])
        
        files = []
        for i in data:
            current_rId = int(i['@Id'].split('Id')[1])
            if (first_slide_id > current_rId) or (current_rId > (first_slide_id+tot_slides-1)):
                files.append(i['@Target'])
        
        target = target + files
        for id in slides:
            slide = "slide"+str(id)+'.xml'
            target.append([i["@Target"] for i in data if slide in i["@Target"] and "http" not in i["@Target"]][0])
            shutil.copy(tmp_path+'/'+file_name+"/ppt/slides/_rels/"+slide+".rels", output_path+'/'+str(render_id)+'/ppt/slides/_rels/')
            add_files(tmp_path+'/'+file_name+"/ppt/slides/_rels/"+slide+".rels",file_name)
    else:
        for i in data:
            if i["@Target"] in target:
                pass
            elif "http" not in i["@Target"]:
                # print("This time: ", i["@Target"])
                target.append(i['@Target'])
                if ".." in i['@Target'] and "xml" in i['@Target']:
                    path = tmp_path+'/'+file_name+"/ppt/"+i['@Target'].split('..')[1].split('/')[1]+"/_rels/"+i['@Target'].split('..')[1].split('/')[2]+".rels"
                    if os.path.exists(path):
                        add_files(path, file_name)
    
    # copy files from tmp dir to output dir
    for i in target:
        if '../' in i:
            if os.path.exists(tmp_path+'/'+file_name+'/ppt/'+i[3:]):
                shutil.copy(tmp_path+'/'+file_name+'/ppt/'+i[3:], output_path+'/'+str(render_id)+'/ppt/'+i[3:].split('/')[0])
        else:
            shutil.copy(tmp_path+'/'+file_name+'/ppt/'+i, output_path+'/'+str(render_id)+'/ppt/'+i.split('/')[0])
        

      

def copy_prep_xml(path):
    """
    copy main relationship and xml file of the deck
    """
    print("COPY_PREP_XML CALLIMG...")
    global tot_slides, first_slide_id, slides
    slides_id = ['rId'+str(first_slide_id+i-1) for i in slides]
    
    # Setting up the paths for xml and rels file
    rels_path = path+'_rels/presentation.xml.rels'
    xml_path = path+'presentation.xml'
    
    # Passing the path of the xml document to enable the parsing process
    # for rels file
    parser = etree.XMLParser(remove_blank_text=True)
    tree = etree.parse(rels_path, parser)
    root = tree.getroot()

    # iterating root
    for relation in root:
        attrib = relation.attrib

        if int(attrib.get('Id').split('Id')[1]) >= first_slide_id and int(attrib.get('Id').split('Id')[1])<(first_slide_id+tot_slides):
            if attrib.get('Id') not in slides_id:
                root.remove(relation)
    tree.write(output_path+'/'+str(render_id)+'/ppt/_rels/presentation.xml.rels', pretty_print=True, xml_declaration=True, encoding='UTF-8')
    
    # Passing the path of the xml document to enable the parsing process
    # for XML file
    parser = etree.XMLParser(remove_blank_text=True)
    tree = etree.parse(xml_path, parser)
    root = tree.getroot()
    for relation in root:
        for ele in relation:
            try:
                rId = int(ele.attrib.values()[-1].split('Id')[-1])
                if rId>=first_slide_id and rId<(first_slide_id+tot_slides):
                    # print("GGG")
                    if 'rId'+str(rId) not in slides_id:
                        relation.remove(ele)
            except:
                pass
        tree.write(output_path+'/'+str(render_id)+'/ppt/presentation.xml', pretty_print=True, xml_declaration=True, encoding='UTF-8')
    print("COMPLETED!!1")


def new (ft) :
    """
    create, move and unzip the empty output deck
    """
    fq_empty = "resources/Empty.pptx"
    # create
    prs = Presentation()
    prs.save(fq_empty)
    # move
    d = ".".join ([ft, "pptx"]) # "output/410.pptx"
    shutil.move (fq_empty, d)
    # unzip
    unzip (d, d.split('.')[0])
    return


def deck_handle(id, msg):
    """
    handle the deck and select files for output deck
    """
    global slides
    file_name, slides = msg['d'], msg['s']
    output_file_loc = f'output/{render_id}'
    # new(output_file_loc)

    # unzip the input deck
    unzip(dir_path+'/presentations/'+file_name+'.pptx', tmp_path+'/'+file_name)
    
    if not os.path.isdir(f'{output_path}/{str(render_id)}'):
        # copy all the necessary files with folder architecture
        shutil.copytree(f'{tmp_path}/{file_name}', f'{output_path}/{str(render_id)}', ignore=ig_d)
        shutil.copytree(f'{tmp_path}/{file_name}/ppt', f'{output_path}/{str(render_id)}/ppt', ignore=ig_f)
    
    path = tmp_path+'/'+file_name+'/ppt/'
    rels_path = tmp_path+'/'+file_name+'/ppt/_rels/presentation.xml.rels'
    
    if slides:
        add_files(rels_path, file_name, slides)
        copy_rel(tmp_path+'/'+file_name+'/ppt', output_path+'/'+str(render_id)+'/ppt')
        copy_mandatory(tmp_path+'/'+file_name+'/ppt/', output_path+'/'+str(render_id)+'/ppt/')
        copy_prep_xml(path)
        # zipdir( output_path+'/41/', file_name)
        
    else:
        o_prs = Presentation(dir_path+'/presentations/'+file_name+'.pptx')
        # o_prs.save('output/'+f'Test_{file_name}.pptx')
    print("TARGET: ", target)
    # remove output/41
    # shutil.rmtree(output_path+'/'+str(render_id))
    
    # print("TARGET", target)
    
    
if __name__ == '__main__':
    
    dir_path = os.path.dirname(os.path.realpath(__file__))
    print("CURRENT_DIR:", dir_path)
    target = []
    
    # load the message
    # file = open('sample_input.json')
    # sample_msg = json.load(file)
    # file.close()
    sample_msg = [41,{'d': 'Onboarding','s':  [2,4,6]}]
    # sample_msg = [41,{'d': 'Presentation1','s':  [1]}]
    # sample_msg = [41,{'d': 'BI Case Studies','s':  [2, 3]}]

    render_id = sample_msg.pop(0)
    output_path = "{}/output".format(dir_path)
    tmp_path = "{}/tmp/{}".format(dir_path, render_id)
    print("TMP_PATH:", tmp_path, '\nOUT_PATH: ', output_path)   

    try:
        os.makedirs(output_path)
        os.makedirs(tmp_path)
    except:
        print("DIR ALREADY EXIST")
    
    # iterating all the messages
    # while sample_msg:
    #     deck_handle(render_id, sample_msg.pop(0))
    # for i in range(1,len(sample_msg)):
    #     deck_handle(render_id, sample_msg[i])
    # zipdir(output_path+'/'+'41', 'Testing')

