import zipfile
import os
import shutil
import xmltodict
import json

from pptx import Presentation
from lxml import etree
from collections import OrderedDict
# from zip_unzip import unzip, zipdir
# from first_deck import copy_mandatory, copy_rel, copy_prep_xml


def unzip(file_path, unzip_path):
    """
    unzip the deck
    """
    with zipfile.ZipFile(file_path, 'r') as zip_ref:
        zip_ref.extractall(unzip_path)
    return


def zipdir(path, file_name):
    """
    zip extracted deck to get output deck
    """
    length = len(path)
    zipf = zipfile.ZipFile('output/'+f'Test_{file_name}.pptx', 'w', zipfile.ZIP_DEFLATED)
    for root, dirs, files in os.walk(path):
        folder = root[length:] # path without "parent"
        for file in files:
            zipf.write(os.path.join(root, file), os.path.join(folder, file))
    zipf.close()
    return


def new (ft) :
    """
    create, move and unzip the empty output deck
    """
    fq_empty = "resources/Empty.pptx"
    # create
    prs = Presentation()
    prs.save(fq_empty)
    # move
    d = ".".join ([ft, "pptx"]) # "output/41.pptx"
    shutil.move (fq_empty, d)
    # unzip
    unzip (d, d.split('.')[0])
    return


# ------- First input deck --------
def make_structure(file):
    """
    creates structure of input deck
    """
    for i in os.walk(f'{tmp_path}/{file}'):
        fld = i[0].split(file)[-1]
        if fld:
            loc = f"{output_file_loc}{fld}"
            if 'ppt' not in loc and (file not in loc):
                shutil.rmtree(f'{output_file_loc}/{fld}')
                shutil.copytree(f'{tmp_path}/{file}/{i[0].split(file)[-1]}', f'{output_file_loc}/{fld}')
            # if not os.path.exists(f'{output_file_loc}/{fld}'):
            #     if 'ppt' not in f'{output_file_loc}/{fld}':
            #         shutil.copytree(f"{tmp_path}/{file}/{i[0].split(f'{file}/')[-1]}", f'{output_file_loc}/{i[0].split(file)[-1]}')
            #     else:
            #         os.makedirs(f'{output_file_loc}/{i[0].split(file)[-1]}')
    return


# ----- both -----
def make_dir(file): # output_file_loc = des
    """
    creates input deck directories which does not exists in the output deck
    """
    # print("777777: ", file, "\nJJJJJ: ", f'{tmp_path}/{file}')
    for i in os.walk(f'{tmp_path}/{file}'):
        fld = i[0].split(file)[-1]
        if fld:
            loc = f"{output_file_loc}{fld}"
            # if os.path.exists(loc) and ('ppt' not in loc) and (file not in loc):
            #     shutil.rmtree(f'{output_file_loc}/{fld}')
            #     shutil.copytree(f"{tmp_path}/{file}/{i[0].split(f'{file}/')[-1]}", f'{output_file_loc}/{i[0].split(file)[-1]}')
            if not os.path.exists(f'{output_file_loc}/{fld}'):
                os.makedirs(f'{output_file_loc}/{i[0].split(file)[-1]}')
                # if 'ppt' not in f'{output_file_loc}/{fld}':
                #     shutil.copytree(f"{tmp_path}/{file}/{i[0].split(f'{file}/')[-1]}", f'{output_file_loc}/{i[0].split(file)[-1]}')
                # else:
                #     os.makedirs(f'{output_file_loc}/{i[0].split(file)[-1]}')
    return


# ----- both -----
def first_slide(path):
    """
    returns the first slide rId
    """
    print("CALLING.. First_slide")
    print("PATH: ", path)
    root, _ = gen_tree(path)

    slide = 'slide1.xml'
    for relation in root:
        attrib = relation.attrib
        if slide in attrib['Target']:
            return int(attrib['Id'].split('Id')[-1])


# ----- both ----- (needs to be modified)
def copy_mandatory(src, des):
    """
    copy mandatory files
    """
    print("COPY MANDATORY CALLING")
    lis = ['slideLayouts', 'theme', 'slideMasters']
    for dir in lis:
        if os.path.exists(f'{des}/{dir}'):
            shutil.rmtree(f'{des}/{dir}')
        shutil.copytree(f'{src}/{dir}', f'{des}/{dir}')
    print("MANDATORY DONE!!")
    
    return


# ----- both -----
def xml_to_dict(path):
    """
    convert xml to dict
    """
    with open(path) as xml_file:
        data_dict = xmltodict.parse(xml_file.read())
        xml_file.close()
    if isinstance(data_dict["Relationships"]["Relationship"], list):
        data = sorted(data_dict["Relationships"]["Relationship"], key=lambda item: int(item['@Id'].split('Id')[1]))
    else:
        data = [data_dict["Relationships"]["Relationship"]]
    return data


# ---- dont know ----
def xml_tag(inp_tag, out_tag):
    """
    returns a dict of tags
    """
    tag_dict = OrderedDict()
    for i in range(len(inp_tag)):
        if inp_tag[i] not in out_tag:
            if i == 0:
                tag_dict[inp_tag[i]] = [0]
            else:
                tag_dict[inp_tag[i]] = [inp_tag[i-1]]
    return tag_dict
    

# ----- both -----
def gen_tree(path):
    """
    pass the path of the xml document to enable the parsing process
    """
    print("CALLING.. Tree")
    parser = etree.XMLParser(remove_blank_text=True)
    tree = etree.parse(path, parser)
    root = tree.getroot()    
    return root, tree


# ----- both -----
def add_files(path, file_name, target_files, slides=None):
    """
    returns a list of files that needs to be modified in output deck
    """
    print("CALLING.. Add_files")
    # global target_files
    data = xml_to_dict(path)
    if slides:
        global sldIds
        # get total slides
        # prs = Presentation(f"{base_path}/presentations/{file_name}.pptx")
        tot_slides = total_slides(f'{input_decks}/{file_name}.pptx')

        first_slide_id = first_slide(path)
        files = []
        for i in data:
            current_rId = int(i['@Id'].split('Id')[1])
            if (first_slide_id > current_rId) or (current_rId > (first_slide_id+tot_slides-1)):
                files.append(i['@Target'])
        
        target_files = target_files + files
        for id in slides:
            slide = f'slide{str(id)}.xml'
            sldIds.append([i["@Id"] for i in data if slide in i["@Target"] and "http" not in i["@Target"]][0])
            target_files.append([i["@Target"] for i in data if slide in i["@Target"] and "http" not in i["@Target"]][0])
            shutil.copy(f'{tmp_path}/{file_name}/ppt/slides/_rels/{slide}.rels', f'{output_path}/{str(render_id)}/ppt/slides/_rels/')
            add_files(f'{tmp_path}/{file_name}/ppt/slides/_rels/{slide}.rels', file_name, target_files)
        # print("UUU1: ", target_files)
    else:
        for i in data:
            if i["@Target"] in target_files:
                pass
            elif "http" not in i["@Target"]:
                # print("This time: ", i["@Target"])
                target_files.append(i['@Target'])
                if ".." in i['@Target'] and "xml" in i['@Target']:
                    path = f"{tmp_path}/{file_name}/ppt/{i['@Target'].split('..')[1].split('/')[1]}/_rels/{i['@Target'].split('..')[1].split('/')[2]}.rels"
                    if os.path.exists(path):
                        # target_files.append(path.split('ppt/')[1])
                        add_files(path, file_name, target_files)
                        
        # print("UUU2: ", target_files)
    # print("UUU: ", target_files)
    return target_files
    

# ------ both ----- (needs to be modified) 
def copy_file(target_files, file_name):
    # copy files from tmp dir to output dir
    for i in target_files:
        if '../' in i:
            if os.path.exists(f'{tmp_path}/{file_name}/ppt/{i[3:]}'):
                shutil.copy(f'{tmp_path}/{file_name}/ppt/{i[3:]}', f"{output_file_loc}/ppt/{i[3:].split('/')[0]}")
            elif os.path.exists(f'{tmp_path}/{file_name}/{i[3:]}'):
                print("IIIII: ", i[3:])
                shutil.copy(f'{tmp_path}/{file_name}/{i[3:]}', f'{output_file_loc}/{i[3:]}')
        else:
            shutil.copy(f'{tmp_path}/{file_name}/ppt/{i}', f'{output_file_loc}/ppt/{i}')
            # shutil.copy(tmp_path+'/'+file_name+'/ppt/'+i, f"{output_file_loc}/ppt/{i.split('/')[0]}")
    return
 

# ----- both ----- (EMPTY)
def select_all():
    """
    if slides is None
    copy all the file content from input to output deck
    """
    pass


# ----- both -----
def total_slides(path):
    """
    returns total number of slides
    """
    print("CALLING.. total_slides")
    prs = Presentation(path)
    tot_slides = len(prs.slides._sldIdLst)
    return tot_slides
    

def copy_prep_xml(src, des, tmp_loc, file_name): # f"{tmp_file_loc}/ppt"
    """
    copy main relationship and xml file of the deck
    """
    print("CALLIMG.. COPY_PREP_XML")
    # global tot_slides, slides
    
    # Setting up the paths for xml and rels file
    rels_path = f'{src}/_rels/presentation.xml.rels'
    xml_path = f'{src}/presentation.xml'
    
    tot_slides = total_slides(f'{input_decks}/{file_name}.pptx')
    print("TOT: ", tot_slides)
    
    # Passing the path of the xml document to enable the parsing process
    # for rels file
    slide_1 = first_slide(rels_path)
    root, tree = gen_tree(rels_path)

    # iterating root
    for relation in root:
        attrib = relation.attrib

        if int(attrib.get('Id').split('Id')[1]) >= slide_1 and int(attrib.get('Id').split('Id')[1])<(slide_1+tot_slides):
            if attrib.get('Id') not in sldIds:
                root.remove(relation)
    tree.write(f'{des}/_rels/presentation.xml.rels', pretty_print=True, xml_declaration=True, encoding='UTF-8')
    
    # Passing the path of the xml document to enable the parsing process
    # for XML file
    root, tree = gen_tree(xml_path)
    print("ROOT: ", root)
    for relation in root:
        for ele in relation:
            # print("ELE: ", ele, ele.attrib, attrib.values())
            try:
                rId = int(ele.attrib.values()[-1].split('Id')[-1])
                # print("RID: ", rId)
                if rId>=slide_1 and rId<(slide_1+tot_slides):
                    # print("GGG")
                    if 'rId'+str(rId) not in sldIds:
                        relation.remove(ele)
            except:
                pass
    tree.write(f'{des}/presentation.xml', pretty_print=True, xml_declaration=True, encoding='UTF-8')
    print("COMPLETED!!1")


# ---- both ----
def del_files(rels_fl, last_fl, path):
    """
    delete extra files
    """
    for i in rels_fl:
        if i[:-5] not in last_fl:
            os.remove(f'{path}/{i}')


# ----- both ------ (modify)
def copy_rel(tmp_loc, out_loc, file_name): # f"{tmp_file_loc}/ppt"
    """
    copy all relelationship files
    """
    src = f'{tmp_loc}/ppt'
    des = f'{out_loc}/ppt'
    for x in os.walk(src):
        folder = x[0].split('ppt')[1]
        # print("FOLDER: ", folder)
        if folder and '_rels' in folder and 'slides' not in folder:
            # print("SRC: ", src+folder, "\nDES: ", des+folder)
            if os.path.exists(des+folder):
                shutil.rmtree(des+folder)
            shutil.copytree(src+folder, des+folder)
    
    # remove empty directories from output dir
    for dir in os.walk(des):
        if not dir[2]:
            if os.path.exists(dir[0]):
                shutil.rmtree(dir[0])
    last = 0
    for i in os.walk(des):
        folder = i[0].split('ppt')[1]
        if folder and '_rels' in folder:
            if len(i[2]) != len(last[2]):
                del_files(i[2], last[2], i[0])
        last = i
            
    
    print("COPY COMPLETED: ")
    copy_mandatory(src, des)
    copy_prep_xml(src, des, tmp_loc, file_name) # f"{tmp_file_loc}/ppt"
    return



def modify(inp_root, out_root, tag_dict, i_tree, o_tree):
    """
    modify content of files
    """
    for i, o in tag_dict.items():
        print("I: ", i, "\nO: ", o)
        # print("TYPE: ", type(i), type(o))
        if o[0] == 0:
            pass
        else:
            subtag1 = o_tree.find(o[0])
            subtag2 = etree.Element(i)
            # for ele in elements:
            #     subtext = etree.SubElement(subtag2, ele)
            # subtext = etree.SubElement(subtag2)
            subtag1.addnext(subtag2)
            
    with open (f'{output_file_loc}/ppt/presentation.xml', 'wb') as f:
        f.write(etree.tostring(out_root, pretty_print = True))
    
    return
    # subtag1 = i_tree.find("subtag1")
    # subtag2 = etree.Element("subtag2", subattrib2="2")
    # subtext = etree.SubElement(subtag2, "subtext")
    # subtext.text = "text2"
    # subtag1.addnext(subtag2)   # Add subtag2 as a following sibling of subtag1

    # print( etree.tostring(tag, pretty_print=True))
    

def pre_xml(file_name): # tmp/41/{file_name}/ppt/presentation.xml
    """
    modify presentation.xml file
    """
    src_xml = f'{tmp_path}/{file_name}/ppt/presentation.xml'
    des_xml = f'{output_file_loc}/ppt/presentation.xml'
    
    inp_root, inp_tree = gen_tree(src_xml)
    out_root, out_tree = gen_tree(des_xml)
    
    print("")
    # for relation in inp_root:
    #     print("FFFF: ", relation.tag, relation.attrib)
    
    inp_tag = [relation.tag for relation in inp_root]
    out_tag = [relation.tag for relation in out_root]
    # print("INP: ", inp_tag, "\nOUT: ", out_tag)
    
    tag_dict = xml_tag(inp_tag, out_tag)
    print("TAG1: ", tag_dict)
    elements = dict()
    if sldIds:
        print("------", sldIds)
        for relation in inp_root:
            if relation.tag in inp_tag:
                for ele in relation:
                    try:
                        print("sss: ", ele.attrib.values()[-1], ele.tag)
                        if 'rId' in ele.attrib.values()[-1]:
                            if ele.attrib.values()[-1] in sldIds:
                                print("REL: ", relation.tag, tag_dict[relation.tag])
                                # print("HHH: ", tag_dict[ele.tag])
                                # print(ele.attrib.values()[-1])
                                # print("TTT: ", [tag_dict[relation.tag], ele.tag, ele.attrib])
                                value = type(tag_dict[relation.tag])
                                print("VALUE: ", value)
                                print("G")
                                tag_dict[relation.tag] = value.extend([ele.tag, ele.attrib])
                                print("H")
                            else:
                                pass
                        else:
                            print("HERE: ", ele.tag)
                            # tag_dict[relation.tag] = tag_dict[relation.tag]
                    # print("ELE: ", ele, ele.attrib, ele.attrib.values(), ele.attrib.values()[-1], type(ele.attrib.values()[-1]))
                        # elif 
                    except:
                        pass
    print("TAG2: ", tag_dict)
    # print("Elements: ", elements)
            # try:
            #     rId = int(ele.attrib.values()[-1].split('Id')[-1])
            #     if rId>=first_slide_id and rId<(first_slide_id+tot_slides):
            #         # print("GGG")
            #         if 'rId'+str(rId) not in slides_id:
            #             relation.remove(ele)
            # except:
            #     pass
                
    modify(inp_root, out_root, tag_dict, inp_tree, out_tree)
    # print("LIS: ", lis)
    return


# ------- first deck -------
def first_deck(path, tmp_file_loc, file_name, slides, target_files):
    """
    handle first deck
    """
    make_structure(file_name)
    target_files = add_files(path, file_name, target_files, slides)
    copy_file(target_files, file_name)
    print("TARGET: ", target_files)
    copy_rel(tmp_file_loc, output_file_loc, file_name)
    
    

def deck_handle(id, msg, deck):
    """
    handle the deck and select files for output deck
    """
    file_name, slides = msg['d'], msg['s']
    target_files = []
    new(output_file_loc)
    
    # unzip the input deck
    unzip(f'{input_decks}/{file_name}.pptx', f'{tmp_path}/{file_name}')
    
    make_dir(file_name)
    
    tmp_file_loc = f'{tmp_path}/{file_name}'
    # print("FILE_LOC: ", file_name)
    prep_xml_path = f'{tmp_file_loc}/ppt/_rels/presentation.xml.rels'
    if deck == 1:
        first_deck(prep_xml_path, tmp_file_loc, file_name, slides, target_files)
    else:
        # print("SSSS: :", target_files, slides, file_name)
        target_files = add_files(prep_xml_path, file_name, target_files, slides)
        
        print("TARGETLLLLLL: ", target_files)
      
        # print("TARGET: ", target_files)

    
    # print("TARGET11: ", a)
    # pre_xml(file_name)
    
    return
    
    
    
    if slides:
        pass
    # add_files()
    else:
        pass
           
    
if __name__ == '__main__':
    
    base_path = os.path.dirname(os.path.realpath(__file__))
    print("CURRENT_DIR:", base_path)
    # target_files = []
    sldIds = []
    
    # load the message
    # file = open('sample_input.json')
    # sample_msg = json.load(file)
    # file.close()
    sample_msg = [41, {'d': 'Onboarding','s':  [2,4,6]}, {'d': 'Presentation1','s':  [1]}]
    # sample_msg = [41, {'d': 'Onboarding','s':  [2,4,6]}]
    # sample_msg = [41, {'d': 'Presentation1','s':  [1]}]
    # sample_msg = [41, {'d': 'BI Case Studies','s':  [2, 3]}]

    render_id = sample_msg.pop(0)
    
    output_path = f'{base_path}/output'
    tmp_path = f'{base_path}/tmp/{render_id}'
    input_decks = f'{base_path}/presentations'
    output_file_loc = f'{output_path}/{render_id}'
    
    print("TMP_PATH:", tmp_path, '\nOUT_PATH: ', output_path)   

    try:
        os.makedirs(tmp_path)
        os.makedirs(output_path)
    except:
        print("DIR ALREADY EXIST")
    
    # iterating all the messages
    deck = 1
    while sample_msg:
        deck_handle(render_id, sample_msg.pop(0), deck)
        deck += 1

    # zipdir(f'{output_file_loc}', "Test")
