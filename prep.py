import os
import shutil
import json
import xmltodict
import natsort
import zipfile
import re
import logging

from pptx import Presentation
from lxml import etree
from collections import OrderedDict
from functools import reduce

logger = logging.getLogger (__name__)

fq_empty = "mob/beagle/test_resources/Empty.pptx"

base_path = os.getcwd()
output_path = f'{base_path}/output/41'
tmp_path = f'{base_path}/tmp/41'

"""
fqfn: ?

"""


def base_name (fqfn) :
    """
    """
    return reduce(lambda x,f : f(x),
                  [fqfn , 
                   os.path.basename, 
                   os.path.splitext]) [0]


def output_path0():
    """
    return the path of the output deck
    """
    path = f'{base_path()}/output/'
    return


def tmp_path0():
    """
    returns the path of the input deck
    """
    path = f'{base_path()}/tmp'
    return path


def base_path0 ():
    """
    returns base path
    """
    return os.getcwd()


def presentation_path(ft):
    """
    returns the path of a specific deck
    """
    path = f'{base_path}/presentations/{ft}.pptx'
    return path


def file_to_work_dir (fqfn) :
    """
    """
    l = ["/", "tmp", base_name (fqfn)]
    return os.path.join (*l) 


def token_to_work_dir (ft) :
    """
    """
    l = ["/", "tmp", ft]
    return os.path.join (*l) 


def token_is_open (ft) :
    """
    """
    fp = token_to_work_dir (ft)
    if not os.path.exists (fp) :
        raise ValueError (f"{ft} is not open")
    return fp


def file_is_open (fn) :
    """
    """
    fp = file_to_work_dir (fn)
    if not os.path.exists (fp) :
        raise ValueError (f"{ft} is not open")
    return fp


def file_is_not_open (fn) :
    """
    """
    fp = file_to_work_dir (fn)
    if os.path.exists (fp) :
        raise ValueError (f"{fn} is already open")
    return fp

###
# def unpack (fqfn) :
def unpack (src, des) :
    """
    unzip the deck
    """
    # fp = file_is_not_open (fqfn)
    with zipfile.ZipFile (src) as zip_ref:
        zip_ref.extractall (des)
    return des


def pack (path, file_name) :
    """
    zip extracted deck to get output deck
    """
    length = len(path)
    d = 'output/'+f'Test_{file_name}.pptx'
    zipf = zipfile.ZipFile (d, 'w', zipfile.ZIP_DEFLATED)
    for root, dirs, files in os.walk(path):
        folder = root[length:] # path without "parent"
        for file in files:
            zipf.write(os.path.join(root, file), os.path.join(folder, file))
    zipf.close()
    logger.debug (f"wrote {d}")
    return d


def close (ft) :
    """
    """
    fp = token_is_open (ft)
    shutil.rmtree (fp)
    logger.debug (f"deleted {ft}")
    return


# Some default functions
def xml_to_dict(path):
    """
    convert xml to dict
    """
    with open(path) as xml_file:
        data_dict = xmltodict.parse(xml_file.read())
        xml_file.close()
    if isinstance(data_dict["Relationships"]["Relationship"], list):
        data = sorted(data_dict["Relationships"]["Relationship"], key=lambda item: int(item['@Id'].split('Id')[1]))
    else:
        data = [data_dict["Relationships"]["Relationship"]]
    return data


def gen_tree(path):
    """
    pass the path of the xml document to enable the parsing process
    """
    # print("CALLING.. Tree")
    parser = etree.XMLParser(remove_blank_text=True)
    tree = etree.parse(path, parser)
    root = tree.getroot()    
    return root, tree


def repair_path(dict_1):
    """
    replacing '\\' with '/'
    """
    dup_dict = dict(dict_1)
    for k,v in dup_dict.items():
        if '\\' in k:
            key = k.replace('\\', '/')
            val = v.replace('\\', '/')
            del dict_1[k]
            dict_1[key] = val
    return dict_1


def create_json(fl, name):
    """
    creates a json files
    """
    obj = json.dumps(fl, indent=4)
    with open(f"new_json/{name}.json", "w") as outfile:
        outfile.write(obj)
    return


def make_dir(file_name): # output_file_loc = des
    """
    creates input deck directories which does not exists in the output deck
    """
    for i in os.walk(f'{tmp_path}/{file_name}'):
        fld = i[0].split(file_name)[-1]
        if fld:
            loc = f"{output_path}{fld}"
            if not os.path.exists(f'{output_path}/{fld}'):
                os.makedirs(f'{output_path}/{fld}')    
    return output_path


def max_rId():
    """
    returns maximum rId
    """
    path = f'{output_path}/ppt/_rels/presentation.xml.rels'
    root, tree = gen_tree(path)

    rIds = []

    for relation in root:
        attrib = relation.attrib
        rId = int(attrib.get('Id').split('Id')[-1])
        rIds.append(rId)
    return {'rId': max(rIds)}


def total_slides(path):
    """
    returns total number of slides
    """
    prs = Presentation(path)
    tt = len(prs.slides._sldIdLst)
    return tt


def first_slide(path):
    """
    returns the first slide rId
    """
    # print("CALLING.. First_slide")
    root, _ = gen_tree(path)

    slide = 'slide1.xml'
    for relation in root:
        attrib = relation.attrib
        if slide in attrib['Target']:
            return int(attrib['Id'].split('Id')[-1])


def new(path):
    """
    create, move and unzip the empty output deck
    """
    fq_empty = "resources/Empty.pptx"
    # create
    prs = Presentation()
    prs.save(fq_empty)
    # move
    d = ".".join ([path, "pptx"]) # "output/41.pptx"
    shutil.move (fq_empty, d)
    # unzip
    unpack (d, d.split('.')[0])
    os.remove(d)
    # m_rId = max_rId()
    return d


# /mnt/input
# /mnt/output

def build_assets(pr, input_deck, assets, slides=None):
    return assets

def build_rels(ft, assets, rels):
    return rels

def build_content_types(ft, assets, content_types):
# def content(tmp_loc, ref_names_dict):
    """
    add content_type in [Contant_Types].xml file
    """

    fl = '[Content_Types].xml'
    inp_path = '/'.join([tmp_loc, fl])
    out_path = '/'.join([output_path, fl])

    cnt_lst = []
    asset_lst = []
    def_att = []
    d = dict()

    root1,tree1 = gen_tree(inp_path)
    root2,tree2 = gen_tree(out_path)

    # get all the extensions belongs to "Default" tag
    for relation in root2:
        if 'Default' in relation.tag:
            def_att.append(relation.attrib['Extension'])
        else:
            break

    for relation in root1:
        if 'Override' in relation.tag:
            attrib = relation.attrib['PartName'][1:]
            try:
                cnt = attrib.split('ppt/')[-1]
                ini = '/ppt/'
            except:
                cnt = attrib
                ini = '/'
            if cnt in ref_names_dict.keys():
                relation.attrib['PartName'] = f'{ini}{ref_names_dict[cnt]}'
                cnt_lst.append(relation)
                # asset_lst.append(relation.attrib['PartName'])
            else:
                cnt_lst.append(relation)
            if relation.attrib['PartName'] not in asset_lst:
                asset_lst.append(relation.attrib['PartName'])
        else:
            attrib = relation.attrib['Extension']
            if attrib not in def_att:
                cnt_lst.append(relation)
                # asset_lst.append(relation.attrib['Extension'])
        # deal with the assest_lst
    # print("AA: ", asset_lst)
    cnt_lst = natsort.natsorted(cnt_lst)
    for ele in cnt_lst:
        prev = tree2.find(ele.tag)
        prev.addnext(ele)

    tree2.write(out_path, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)

    unq_attr = []
    for relation in root2:
        if 'Override' in relation.tag:
            if relation.attrib['PartName'] not in unq_attr:
                unq_attr.append(relation.attrib['PartName'])
            else:
                root2.remove(relation)
    tree2.write(out_path, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)

    return content_types

def build_properties(ft, assets, properties):
    return properties

def apply_assets(ft, assets):
    return assets

def apply_rels(ft, rels):
    return rels

def apply_content_types(ft, content_types):
    return content_types

def apply_properties(ft, properties):
    return properties


def process_message(msg, ctx):
    input_deck = msg["d"]
    
    ft = unpack (input_deck)

    slides = msg.get("s", None)
    if not slides:
        pt = presentation_path(input_deck)
        slides = total_slides(pt)

    assets_name = OrderedDict()
    assets_cnt = OrderedDict()
    # assets = OrderedDict()
    # rels = OrderedDict()
    # content_types = OrderedDict()
    # properties = []

    assets = []
    pr = f'{tmp_path}/{input_deck}/ppt/_rels/presentation.xml.rels'
    assets = build_assets(pr, input_deck, assets, slides=None)
    
    for s in slides:
        assets = build_assets(s, ft, assets)
        rels = build_rels(ft, assets, rels)
    content_types = build_content_types(ft, assets, content_types)
    properties = build_properties(ft, assets, properties)

    ctx["assets_name"] = {**ctx["assets"], assets}
    ctx["assets_cnt"] = {**ctx["assets"], assets}
    
    # ctx["assets"] = {**ctx["assets"], assets}
    # ctx["rels"] = {**ctx["rels"], rels)
    # ctx["content_types"] = {ctx["content_types"], rels}
    # ctx["properties"] = ctx.get("properties") + properties

    return ctx

def write_output(ctx, output_deck):
    ft = new(output_deck)
    apply_assets(ft, ctx["assets"])
    apply_rels(ft, ctx["rels"])
    apply_content_types(ft, ctx["content_types"])
    apply_properties(ft, ctx["properties"])
    pack(ft, file=output_deck)
    close(ft)
    return output_deck

def render_deck_effect(msgs, fqfn):
    """
    process msgs and write output deck to fqfn
    """
    ctx = reduce(process_message, msgs, {})   
    write_output(ctx, fqfn)
    return fqfn

# continue
